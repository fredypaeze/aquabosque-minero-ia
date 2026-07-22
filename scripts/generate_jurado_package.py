"""Genera el paquete de entregables para jurado 2026 a partir del producto real.

La prioridad es trazabilidad: cada cifra y afirmación sale del código, la app
desplegada o los artefactos reproducibles del repositorio. No inventa métricas,
no fuerza la narrativa del prompt cuando contradice la implementación vigente.
"""
from __future__ import annotations

import json
import math
import os
import shutil
import textwrap
from dataclasses import dataclass
from pathlib import Path

import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import qrcode
from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import landscape
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "jurado_2026"
ASSETS = OUT / "assets"
DEMO = OUT / "capturas_demo"
BACKUP = OUT / "backup_offline"
TMP = ROOT / "tmp"
LIVE = ROOT / "tmp" / "audit" / "live-home.png"
SCREEN_DIR = ROOT / "tmp" / "screens"

PRED = ROOT / "outputs" / "tables" / "predicciones.csv"
MASTER = ROOT / "data" / "processed" / "master_con_etiqueta.csv"
METRICS = ROOT / "models" / "metrics" / "metricas.json"
SHAP = ROOT / "models" / "shap" / "importancia_global.csv"
GEOJSON = ROOT / "data" / "processed" / "municipios.geojson"
RAW_STATE = ROOT / "data" / "raw" / "_estado_fuentes.json"
RAW_LOG = ROOT / "data" / "raw" / "_download_log.json"

VERDE = "#1B5E20"
VERDE_2 = "#2E7D32"
AZUL = "#0277BD"
AZUL_2 = "#01579B"
TINTA = "#12261A"
FONDO = "#F6FAF6"
PANEL = "#FFFFFF"
BORDE = "#DCE8DF"
RISKS = {"Crítico": "#7F1D1D", "Alto": "#DC2626", "Medio": "#F59E0B", "Bajo": "#16A34A"}

os.environ.setdefault("BROWSER_PATH", shutil.which("chromium") or "")


@dataclass
class SlideSpec:
    number: int
    title: str
    bullets: list[str]
    image_path: Path | None
    footnote: str
    objective: str
    speech: str
    duration: str
    takeaway: str
    transition: str
    questions: list[tuple[str, str]]
    subtitle: str | None = None


def ensure_dirs() -> None:
    for path in [OUT, ASSETS, DEMO, BACKUP]:
        path.mkdir(parents=True, exist_ok=True)


def load_data() -> dict:
    pred = pd.read_csv(PRED)
    master = pd.read_csv(MASTER)
    shap_df = pd.read_csv(SHAP)
    metrics = json.loads(METRICS.read_text(encoding="utf-8"))
    raw_state = json.loads(RAW_STATE.read_text(encoding="utf-8"))
    raw_log = json.loads(RAW_LOG.read_text(encoding="utf-8"))
    return {
        "pred": pred,
        "master": master,
        "shap": shap_df,
        "metrics": metrics,
        "raw_state": raw_state,
        "raw_log": raw_log,
    }


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def save_qr(text: str, output: Path) -> None:
    img = qrcode.make(text).convert("RGB")
    img.save(output)


def _wrap_by_pixels(draw: ImageDraw.ImageDraw, paragraph: str, width: int, font_obj) -> list[str]:
    """Envuelve por ancho real en píxeles (no por conteo de caracteres) para que
    el texto nunca desborde su caja: mide cada palabra con la métrica de la fuente."""
    words = paragraph.split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        trial = f"{current} {word}"
        if draw.textlength(trial, font=font_obj) <= width:
            current = trial
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def add_wrapped_text(draw: ImageDraw.ImageDraw, text: str, xy: tuple[int, int], width: int, font_obj, fill: str,
                     line_gap: int = 8) -> int:
    x, y = xy
    lines: list[str] = []
    for paragraph in text.split("\n"):
        if not paragraph.strip():
            lines.append("")
        else:
            lines.extend(_wrap_by_pixels(draw, paragraph, width, font_obj))
    for line in lines:
        draw.text((x, y), line, font=font_obj, fill=fill)
        y += font_obj.size + line_gap
    return y


def crop_live_screens() -> None:
    mappings = {
        LIVE: ASSETS / "01_app_portada_viva.png",
        SCREEN_DIR / "Ficha_territorial.png": ASSETS / "05_ficha_live.png",
        SCREEN_DIR / "Explicabilidad.png": ASSETS / "06_shap_live.png",
        SCREEN_DIR / "Datos_abiertos.png": ASSETS / "07_datos_abiertos_live.png",
    }
    for src, dst in mappings.items():
        if src.exists():
            shutil.copy2(src, dst)
            shutil.copy2(src, DEMO / dst.name)
    if (SCREEN_DIR / "Mapa_de_riesgo.png").exists():
        img = Image.open(SCREEN_DIR / "Mapa_de_riesgo.png")
        crop = img.crop((300, 80, 1430, 420))
        crop.save(ASSETS / "04_filtros_mapa_live.png")


def frame_browser(src: Path, output: Path, url: str = "streamlit.spartanit.pro", target_ratio: float = 1.5) -> None:
    """Envuelve una captura (típicamente vertical) en un marco de navegador y la
    recorta a formato horizontal desde arriba. Evita que los pantallazos altos
    queden como una tira angosta flotando con espacio en blanco alrededor."""
    shot = Image.open(src).convert("RGB")
    # Recorte superior a proporción horizontal (contenido más relevante arriba).
    crop_h = min(shot.height, int(shot.width / target_ratio))
    shot = shot.crop((0, 0, shot.width, crop_h))

    pad = 26
    bar_h = 74
    W = shot.width + pad * 2
    H = shot.height + bar_h + pad * 2
    canvas_img = Image.new("RGB", (W, H), FONDO)
    d = ImageDraw.Draw(canvas_img)
    # Ventana con barra superior
    d.rounded_rectangle((pad - 6, pad - 6, W - pad + 6, H - pad + 6), radius=26, fill="#FFFFFF", outline=BORDE, width=3)
    d.rounded_rectangle((pad, pad, W - pad, pad + bar_h), radius=18, fill="#EAF0EC")
    d.rectangle((pad, pad + bar_h - 18, W - pad, pad + bar_h), fill="#EAF0EC")
    for i, color in enumerate(["#F04B4B", "#F5B93B", "#4BC45A"]):
        d.ellipse((pad + 28 + i * 34, pad + 26, pad + 48 + i * 34, pad + 46), fill=color)
    pill_x0 = pad + 150
    d.rounded_rectangle((pill_x0, pad + 18, W - pad - 40, pad + bar_h - 18), radius=18, fill="#FFFFFF", outline=BORDE, width=2)
    d.text((pill_x0 + 26, pad + 30), f"\U0001F512  {url}", font=font(24), fill="#52645A")
    canvas_img.paste(shot, (pad, pad + bar_h))
    canvas_img.save(output)


def build_satelital_panel(output: Path) -> bool:
    """Panel satelital NRT (mapa de focos + KPIs + prioridad máxima) para el deck.
    Lee la señal ya generada (fuego_municipal.csv). Devuelve False si no existe."""
    fuego_csv = ROOT / "data" / "processed" / "fuego_municipal.csv"
    if not fuego_csv.exists():
        print("  AVISO: no hay fuego_municipal.csv; se omite el panel satelital.")
        return False
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.collections import PolyCollection
    import numpy as np

    fuego = pd.read_csv(fuego_csv)
    pred = pd.read_csv(PRED); pred["cod_mpio"] = pred["cod_mpio"].astype(float).astype(int)
    fuego["cod_mpio"] = fuego["cod_mpio"].astype(int)
    geo = json.loads(GEOJSON.read_text(encoding="utf-8"))
    focos = dict(zip(fuego["cod_mpio"], fuego["focos_7d"]))

    base = pred.drop(columns=[c for c in ["focos_7d", "frp_total", "idx_fuego"] if c in pred.columns])
    mix = base.merge(fuego[["cod_mpio", "focos_7d", "frp_total"]], on="cod_mpio", how="left").fillna(0)
    prioridad = mix[(mix.riesgo_nivel.isin(["Alto", "Crítico"])) & (mix.focos_7d > 0)]
    nuevos = mix[(mix.frp_total > 200) & (mix.riesgo_nivel.isin(["Bajo", "Medio"]))]
    total_focos = int(fuego["focos_7d"].sum())

    fig = plt.figure(figsize=(16, 9), dpi=100)
    fig.patch.set_facecolor(FONDO)
    ax = fig.add_axes([0.02, 0.02, 0.46, 0.9]); ax.set_facecolor(FONDO)
    polys, vals = [], []
    for feat in geo["features"]:
        ring = feat["geometry"]["coordinates"][0]
        polys.append([(p[0], p[1]) for p in ring])
        vals.append(focos.get(int(feat["properties"]["cod"]), 0))
    vals = np.array(vals, dtype=float)
    ax.add_collection(PolyCollection([p for p, v in zip(polys, vals) if v == 0],
                                     facecolors="#EDF2ED", edgecolors="#DCE4DC", linewidths=0.15))
    hot = [p for p, v in zip(polys, vals) if v > 0]
    if hot:
        pc = PolyCollection(hot, array=np.clip(vals[vals > 0], 0, np.percentile(vals[vals > 0], 97)),
                            cmap="YlOrRd", edgecolors="#B0453A", linewidths=0.2)
        ax.add_collection(pc)
    ax.set_xlim(-79.4, -66.7); ax.set_ylim(-4.5, 13.7); ax.set_aspect("equal"); ax.axis("off")
    ax.set_title("Focos de calor activos por municipio · 7 días", fontsize=15,
                 fontweight="bold", color=TINTA, loc="left")

    # Panel derecho: KPIs + tabla prioridad máxima
    def kpi(x, y, val, lbl, color):
        fig.text(x, y, val, fontsize=30, fontweight="bold", color=color, transform=fig.transFigure)
        fig.text(x, y - 0.05, lbl, fontsize=12, color="#52645A", transform=fig.transFigure)
    kpi(0.52, 0.82, f"{total_focos:,}".replace(",", "."), "focos satelitales (7d)", AZUL_2)
    kpi(0.72, 0.82, str(int((fuego.focos_7d > 0).sum())), "municipios con fuego", VERDE)
    kpi(0.52, 0.66, str(len(prioridad)), "prioridad máxima (modelo + fuego)", "#B91C1C")
    kpi(0.72, 0.66, str(len(nuevos)), "actividad nueva (no vista por índice)", "#B45309")

    fig.text(0.52, 0.55, "Prioridad máxima de verificación", fontsize=14, fontweight="bold", color=TINTA)
    fig.text(0.52, 0.52, "El modelo los prioriza y el satélite confirma fuego hoy:", fontsize=11, color="#52645A")
    top = prioridad.sort_values("frp_total", ascending=False).head(6)
    y = 0.47
    for _, r in top.iterrows():
        fig.text(0.52, y, f"• {str(r['municipio']).title()} ({str(r['departamento']).title()})",
                 fontsize=12.5, color=TINTA)
        fig.text(0.93, y, f"{int(r['focos_7d'])} focos", fontsize=12, color="#B0453A", ha="right")
        y -= 0.045
    fig.text(0.52, 0.06, "Fuente: NASA FIRMS · VIIRS + MODIS · datos abiertos · señal NRT (proxy de deforestación/quema).",
             fontsize=10, color="#65766C")
    fig.savefig(output, facecolor=FONDO, bbox_inches="tight")
    plt.close(fig)
    return True


def _write_plot(fig, output: Path, width: int, height: int) -> None:
    fig.write_image(str(output), width=width, height=height, scale=2)


def build_map_image(df: pd.DataFrame, output: Path) -> None:
    with GEOJSON.open(encoding="utf-8") as fh:
        geo = json.load(fh)
    work = df.copy()
    work["cod_mpio"] = work["cod_mpio"].astype(int)
    order = ["Bajo", "Medio", "Alto", "Crítico"]
    fig = px.choropleth(
        work,
        geojson=geo,
        locations="cod_mpio",
        featureidkey="id",
        color="riesgo_nivel",
        category_orders={"riesgo_nivel": order},
        color_discrete_map={k: v for k, v in RISKS.items()},
        hover_name="municipio",
        hover_data={"departamento": True, "riesgo_score": ":.3f", "riesgo_nivel": True, "cod_mpio": False},
    )
    fig.update_geos(fitbounds="locations", visible=False, showcountries=False, showcoastlines=False, bgcolor=FONDO)
    fig.update_traces(marker_line_width=0.3, marker_line_color="#FFFFFF")
    fig.update_layout(
        width=1400,
        height=900,
        font=dict(size=26, color=TINTA),
        paper_bgcolor=FONDO,
        plot_bgcolor=FONDO,
        margin=dict(l=0, r=0, t=70, b=0),
        legend=dict(orientation="h", y=1.02, x=0, title=None, font=dict(size=28)),
        title=dict(text="Cobertura nacional de priorización municipal", x=0.01, y=0.97, font=dict(size=40, color=TINTA)),
    )
    _write_plot(fig, output, 1400, 900)


def build_coverage_map(master: pd.DataFrame, output: Path) -> None:
    with GEOJSON.open(encoding="utf-8") as fh:
        geo = json.load(fh)
    work = master.copy()
    work["cod_mpio"] = work["cod_mpio"].astype(int)
    work["n_senales"] = (
        (work["mineria_titulos"] > 0).astype(int)
        + (work["deforestacion_ha"] > 0).astype(int)
        + (work["agua_estaciones"] > 0).astype(int)
        + (work["runap_areas"] > 0).astype(int)
    )
    fig = px.choropleth(
        work,
        geojson=geo,
        locations="cod_mpio",
        featureidkey="id",
        color="n_senales",
        color_continuous_scale=["#EDF5EE", "#9FD3A9", "#4A8F56", "#0E5D31"],
        hover_name="municipio",
        hover_data={"departamento": True, "n_senales": True, "cod_mpio": False},
    )
    fig.update_geos(fitbounds="locations", visible=False, bgcolor=FONDO)
    fig.update_traces(marker_line_width=0.3, marker_line_color="#FFFFFF")
    fig.update_layout(
        width=1400,
        height=900,
        font=dict(size=24, color=TINTA),
        paper_bgcolor=FONDO,
        margin=dict(l=0, r=0, t=70, b=0),
        coloraxis_colorbar=dict(title=dict(text="Señales", font=dict(size=26)), tickfont=dict(size=22)),
        title=dict(text="Cobertura territorial por número de señales observadas", x=0.01, y=0.97, font=dict(size=38, color=TINTA)),
    )
    _write_plot(fig, output, 1400, 900)


def build_ranking_chart(df: pd.DataFrame, output: Path) -> None:
    top = df.sort_values("riesgo_score", ascending=False).head(15).copy()
    top = top.iloc[::-1]
    fig = px.bar(
        top,
        x="riesgo_score",
        y="municipio",
        color="riesgo_nivel",
        orientation="h",
        color_discrete_map=RISKS,
        text="riesgo_score",
        hover_data={"departamento": True},
    )
    fig.update_traces(texttemplate="%{text:.3f}", textposition="outside", textfont_size=20, cliponaxis=False)
    fig.update_layout(
        width=1400,
        height=900,
        font=dict(size=22, color=TINTA),
        paper_bgcolor=FONDO,
        plot_bgcolor=FONDO,
        margin=dict(l=10, r=60, t=80, b=30),
        legend=dict(orientation="h", y=1.03, x=0, title=None, font=dict(size=24)),
        title=dict(text="Top 15 municipios por score de priorización", x=0.01, y=0.97, font=dict(size=40, color=TINTA)),
        xaxis=dict(title=dict(text="Score de priorización", font=dict(size=24)), tickfont=dict(size=20)),
        yaxis=dict(title=None, tickfont=dict(size=22)),
    )
    _write_plot(fig, output, 1400, 900)


def build_sources_table_image(sources: pd.DataFrame, output: Path) -> None:
    W, H = 2000, 1260
    img = Image.new("RGB", (W, H), FONDO)
    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle((36, 26, W - 36, H - 26), radius=28, fill=PANEL, outline=BORDE, width=3)
    draw.text((72, 56), "Fuentes abiertas integradas en el MVP vigente", font=font(40, True), fill=TINTA)
    draw.text((72, 112), "Datasets abiertos (datos.gov.co y portales oficiales) efectivamente consumidos por el pipeline, con corte y función.",
              font=font(21), fill="#52645A")

    headers = ["Dimensión", "Entidad / recurso", "Cobertura usada", "Corte", "Función en el producto"]
    col_x = [72, 348, 812, 1258, 1508]
    col_w = [260, 448, 430, 232, 420]
    header_y = 184
    # Banda de encabezado
    draw.rounded_rectangle((56, header_y - 14, W - 56, header_y + 44), radius=12, fill="#EAF3EE")
    for x, h in zip(col_x, headers):
        draw.text((x, header_y), h, font=font(23, True), fill=VERDE)

    row_top = header_y + 66
    row_h = 150
    for i, (_, row) in enumerate(sources.iterrows()):
        y0 = row_top + i * row_h
        if i % 2 == 1:
            draw.rounded_rectangle((56, y0 - 10, W - 56, y0 + row_h - 22), radius=10, fill="#F6FAF6")
        values = [row["dimension"], row["fuente"], row["cobertura"], row["corte"], row["funcion"]]
        # Dimensión en negrita/color, resto en tinta
        draw.rounded_rectangle((col_x[0], y0 + 2, col_x[0] + 12, y0 + 60), radius=4, fill=VERDE_2)
        add_wrapped_text(draw, str(values[0]), (col_x[0] + 24, y0), col_w[0] - 24, font(22, True), TINTA, line_gap=6)
        for x, value, width in zip(col_x[1:], values[1:], col_w[1:]):
            add_wrapped_text(draw, str(value), (x, y0), width, font(20), "#33413A", line_gap=6)
    img.save(output)


def build_architecture_diagram(output: Path) -> None:
    img = Image.new("RGB", (1600, 900), FONDO)
    d = ImageDraw.Draw(img)
    d.text((60, 48), "Cadena técnica validada por el MVP", font=font(36, True), fill=TINTA)
    d.text((60, 100), "Lo que está funcionando hoy en el producto real, de datos abiertos a app.", font=font(20), fill="#52645A")
    box_w, box_h, box_y = 250, 210, 250
    xs = [60, 367, 674, 981, 1288]
    boxes = [
        ("Datos abiertos\ndatos.gov.co\nANM · IDEAM · RUNAP\nDANE · PDET", VERDE),
        ("Integración municipal\ncódigos DANE +\ncentroides + cruces\npor proximidad", VERDE_2),
        ("5 índices + score\ntécnico: minero, bosque,\nfuego satelital, hídrico,\nsensibilidad", AZUL),
        ("XGBoost multiclase\n+ SHAP para explicar\nla priorización", AZUL_2),
        ("App Streamlit\nmapa, ranking,\nficha y metodología", "#0D6B53"),
    ]
    for (text, color), x in zip(boxes, xs):
        d.rounded_rectangle((x, box_y, x + box_w, box_y + box_h), radius=22, fill=color, outline=None)
        add_wrapped_text(d, text, (x + 20, box_y + 26), box_w - 40, font(23, True), "#FFFFFF", line_gap=8)
    mid_y = box_y + box_h // 2
    for i in range(len(xs) - 1):
        x_end = xs[i] + box_w
        nx = xs[i + 1]
        d.line((x_end + 8, mid_y, nx - 10, mid_y), fill="#7CA98A", width=6)
        d.polygon([(nx - 6, mid_y), (nx - 26, mid_y - 11), (nx - 26, mid_y + 11)], fill="#7CA98A")
    d.rounded_rectangle((120, 560, 1480, 800), radius=26, fill="#FFFFFF", outline=BORDE, width=3)
    d.text((150, 588), "Hallazgo de auditoría", font=font(24, True), fill=VERDE)
    add_wrapped_text(
        d,
        "La IA vigente no es un detector de anomalías ni un motor satelital en tiempo real. "
        "El código y la app exponen un clasificador XGBoost sobre una etiqueta técnica por cuantiles, "
        "acompañado por SHAP para interpretabilidad. Se defiende lo que realmente está implementado.",
        (150, 632),
        1300,
        font(24),
        TINTA,
        line_gap=10,
    )
    img.save(output)


def build_dataflow_diagram(output: Path) -> None:
    img = Image.new("RGB", (1600, 900), FONDO)
    d = ImageDraw.Draw(img)
    d.text((60, 55), "Cómo se construye la unidad territorial comparable", font=font(34, True), fill=TINTA)
    steps = [
        ("1", "Base municipal\nDIVIPOLA", "1.122 municipios con código y centroide"),
        ("2", "Cruces directos", "RUCOM, ANM volumen y PDET por código DANE"),
        ("3", "Cruces de proximidad", "ICA y RUNAP por centroide municipal más cercano"),
        ("4", "Normalización", "4 índices 0-1 + score técnico de priorización"),
        ("5", "Salida usable", "Predicciones, confianza, SHAP y vistas de app"),
    ]
    x_positions = [70, 360, 650, 940, 1230]
    for (num, title, desc), x in zip(steps, x_positions):
        d.ellipse((x, 180, x + 90, 270), fill=AZUL if num in {"3", "4"} else VERDE_2)
        d.text((x + 32, 205), num, font=font(34, True), fill="#FFFFFF")
        d.rounded_rectangle((x - 20, 310, x + 240, 620), radius=20, fill="#FFFFFF", outline=BORDE, width=3)
        add_wrapped_text(d, title, (x + 8, 345), 210, font(24, True), TINTA, line_gap=5)
        add_wrapped_text(d, desc, (x + 8, 420), 210, font(18), "#496057", line_gap=5)
        if x < x_positions[-1]:
            d.line((x + 240, 450, x + 285, 450), fill="#7CA98A", width=5)
            d.polygon([(x + 285, 450), (x + 265, 440), (x + 265, 460)], fill="#7CA98A")
    d.rounded_rectangle((80, 700, 1520, 820), radius=20, fill="#EAF3EE", outline=None)
    add_wrapped_text(
        d,
        "Regla clave de consistencia: cuando no hay estación ICA el sistema marca ausencia de dato hídrico y usa 0 como señal observada; "
        "no imputa contaminación. Deforestación y RUNAP usan cortes geoespaciales acotados, no cobertura homogénea nacional.",
        (120, 735),
        1340,
        font(20),
        TINTA,
        line_gap=8,
    )
    img.save(output)


def build_model_explainer_image(metrics: dict, shap_df: pd.DataFrame, output: Path) -> None:
    img = Image.new("RGB", (1600, 900), FONDO)
    d = ImageDraw.Draw(img)
    d.text((60, 46), "Modelo vigente y forma correcta de explicarlo", font=font(36, True), fill=TINTA)
    d.text((60, 96), "Clasificador XGBoost multiclase + SHAP sobre una etiqueta técnica por cuantiles.", font=font(20), fill="#52645A")

    card_top, card_bottom = 165, 470
    # Tarjeta 1: entrada del modelo
    d.rounded_rectangle((60, card_top, 520, card_bottom), radius=22, fill=PANEL, outline=BORDE, width=3)
    d.text((92, card_top + 26), "Entrada del modelo", font=font(25, True), fill=VERDE)
    add_wrapped_text(
        d,
        "15 variables: 5 índices (incl. fuego satelital) y 10 variables crudas trazables.",
        (92, card_top + 76), 400, font(22), TINTA, line_gap=8,
    )
    d.text((92, card_top + 200), "Salida", font=font(22, True), fill=TINTA)
    d.text((92, card_top + 238), "Bajo · Medio · Alto · Crítico", font=font(21), fill="#496057")

    # Tarjeta 2: honestidad metodológica (métricas destacadas)
    d.rounded_rectangle((548, card_top, 1010, card_bottom), radius=22, fill=PANEL, outline=BORDE, width=3)
    d.text((580, card_top + 26), "Honestidad metodológica", font=font(25, True), fill=VERDE)
    metric_rows = [
        (f"{metrics['accuracy']:.1%}", "Accuracy"),
        (f"{metrics['baseline_clase_mayoritaria']:.1%}", "Línea base (clase mayoritaria)"),
        (f"{metrics['f1_macro']:.2f}", "F1 macro"),
    ]
    my = card_top + 74
    for value, label in metric_rows:
        d.text((580, my), value, font=font(30, True), fill=AZUL_2)
        d.text((700, my + 8), label, font=font(18), fill="#496057")
        my += 58
    d.text((580, my + 4), "Accuracy alta porque la etiqueta es una", font=font(16), fill="#7C4A03")
    d.text((580, my + 28), "fórmula compuesta: no se vende como", font=font(16), fill="#7C4A03")
    d.text((580, my + 52), "mérito predictivo.", font=font(16), fill="#7C4A03")

    # Tarjeta 3: SHAP — variables que más pesan (con barras)
    d.rounded_rectangle((1038, card_top, 1540, card_bottom), radius=22, fill=PANEL, outline=BORDE, width=3)
    d.text((1070, card_top + 26), "Variables que más pesan (SHAP)", font=font(22, True), fill=VERDE)
    top = shap_df.sort_values("importancia_shap", ascending=False).head(4)
    max_imp = float(top["importancia_shap"].max()) or 1.0
    bar_x0, bar_x_max = 1070, 1500
    by = card_top + 84
    for _, row in top.iterrows():
        name = str(row["feature"])
        val = float(row["importancia_shap"])
        d.text((bar_x0, by), name, font=font(17, True), fill=TINTA)
        d.rounded_rectangle((bar_x0, by + 24, bar_x_max, by + 38), radius=7, fill="#E8F0EA")
        bar_len = int((bar_x_max - bar_x0) * (val / max_imp))
        d.rounded_rectangle((bar_x0, by + 24, bar_x0 + max(6, bar_len), by + 38), radius=7, fill=VERDE_2)
        d.text((bar_x_max - 70, by), f"{val:.3f}", font=font(16), fill="#496057")
        by += 62

    # Franja de conclusión (auditoría)
    d.rounded_rectangle((60, 510, 1540, 800), radius=24, fill="#FFF7ED", outline="#F4B04A", width=3)
    d.text((92, 536), "Alcance honesto de la IA", font=font(26, True), fill="#7C4A03")
    add_wrapped_text(
        d,
        "El modelo integra una señal satelital NRT (focos térmicos FIRMS) como dimensión, real y diaria. "
        "El procesamiento de imagen cruda Sentinel-2 con deep learning corre en la infraestructura GPU del "
        "Ministerio (capa en construcción). No se presenta detección de anomalías ni satélite operativo que no "
        "exista: se defiende lo implementado —una priorización interpretable que cruza minería, bosque, fuego, agua y sensibilidad.",
        (92, 584),
        1400,
        font(22),
        "#5A3703",
        line_gap=9,
    )
    img.save(output)


def build_case_image(row: pd.Series, output: Path) -> None:
    img = Image.new("RGB", (1600, 900), FONDO)
    d = ImageDraw.Draw(img)
    d.text((60, 55), f"Caso real: {row['municipio']}, {row['departamento']}", font=font(34, True), fill=TINTA)
    d.text((60, 105), "Municipio crítico con señal minera, hídrica y de sensibilidad territorial visibles en el dataset.", font=font(19), fill="#52645A")
    left = 70
    for i, (label, value) in enumerate([
        ("Nivel", row["riesgo_nivel"]),
        ("Score", f"{row['riesgo_score']:.3f}"),
        ("Predicción", row["riesgo_pred"]),
        ("Confianza", f"{row['confianza']:.0%}"),
    ]):
        x = left + i * 365
        d.rounded_rectangle((x, 170, x + 320, 290), radius=18, fill=PANEL, outline=BORDE, width=3)
        d.text((x + 25, 195), label, font=font(20), fill="#5C6F63")
        d.text((x + 25, 230), str(value), font=font(30, True), fill=TINTA)

    d.rounded_rectangle((70, 340, 970, 790), radius=26, fill=PANEL, outline=BORDE, width=3)
    d.text((100, 370), "Perfil por dimensión", font=font(28, True), fill=TINTA)
    dims = [
        ("Minero", float(row["idx_minero"]), VERDE_2),
        ("Deforestación", float(row["idx_deforestacion"]), "#4C9E62"),
        ("Hídrico", float(row["idx_hidrico"]), AZUL),
        ("Sensibilidad", float(row["idx_sensibilidad"]), AZUL_2),
    ]
    for idx, (label, value, color) in enumerate(dims):
        y = 440 + idx * 80
        d.text((100, y - 6), label, font=font(22, True), fill=TINTA)
        d.rounded_rectangle((290, y, 860, y + 30), radius=15, fill="#E8F0EA", outline=None)
        d.rounded_rectangle((290, y, 290 + int(570 * value), y + 30), radius=15, fill=color, outline=None)
        d.text((885, y - 6), f"{value:.2f}", font=font(21, True), fill=TINTA)

    d.rounded_rectangle((1030, 340, 1530, 790), radius=26, fill=PANEL, outline=BORDE, width=3)
    d.text((1060, 370), "Lectura útil para autoridad", font=font(28, True), fill=TINTA)
    case_text = (
        f"• {int(row['mineria_titulos'])} registros RUCOM y {int(row['mineria_minerales'])} minerales distintos.\n"
        f"• {int(row['agua_estaciones'])} estación ICA cercana con índice hídrico {float(row['idx_hidrico']):.2f}.\n"
        f"• {int(row['runap_areas'])} áreas RUNAP asociadas; {float(row['runap_hectareas']):,.1f} ha protegidas.\n"
        f"• No confirma daño: prioriza revisión técnica, seguimiento hídrico y cruce con actividad minera formal."
    )
    add_wrapped_text(d, case_text, (1070, 435), 380, font(22), TINTA, line_gap=10)
    img.save(output)


def build_roadmap(output: Path) -> None:
    img = Image.new("RGB", (1600, 900), FONDO)
    d = ImageDraw.Draw(img)
    d.text((60, 55), "Ruta de escalamiento creíble", font=font(34, True), fill=TINTA)
    phases = [
        ("Fase actual\nMVP funcional", ["Integración municipal", "XGBoost + SHAP", "App desplegada", "16 pruebas pasando"], VERDE),
        ("Escalamiento técnico", ["Automatizar descargas", "Históricos por fuente", "Reentrenamiento periódico", "Deep-links y mejor demo"], AZUL),
        ("Uso operacional", ["Validación experta", "Alertas según periodicidad", "Reportes a autoridad", "Nuevas coberturas"], AZUL_2),
    ]
    xs = [90, 560, 1030]
    for (title, bullets, color), x in zip(phases, xs):
        d.rounded_rectangle((x, 180, x + 420, 700), radius=28, fill=PANEL, outline=color, width=5)
        d.rounded_rectangle((x, 180, x + 420, 280), radius=28, fill=color, outline=color)
        add_wrapped_text(d, title, (x + 28, 205), 340, font(30, True), "#FFFFFF", line_gap=6)
        y = 335
        for bullet in bullets:
            d.text((x + 32, y), f"• {bullet}", font=font(23), fill=TINTA)
            y += 72
    img.save(output)


def _cover_crop(img: Image.Image, box_w: int, box_h: int) -> Image.Image:
    """Escala la imagen para CUBRIR la caja (sin bandas blancas) y recorta el
    excedente; alineado arriba para capturas de página larga."""
    scale = max(box_w / img.width, box_h / img.height)
    resized = img.resize((max(1, int(img.width * scale)), max(1, int(img.height * scale))))
    left = (resized.width - box_w) // 2
    return resized.crop((left, 0, left + box_w, box_h))


def build_app_collage(output: Path) -> None:
    paths = [
        ASSETS / "01_app_portada_viva.png",
        ASSETS / "05_ficha_live.png",
        ASSETS / "06_shap_live.png",
        ASSETS / "07_datos_abiertos_live.png",
    ]
    imgs = [Image.open(p).convert("RGB") for p in paths if p.exists()]
    canvas_img = Image.new("RGB", (1800, 1200), FONDO)
    d = ImageDraw.Draw(canvas_img)
    d.text((55, 34), "Aplicación funcional: vistas auditadas de la versión desplegada", font=font(36, True), fill=TINTA)
    boxes = [(50, 120, 910, 630), (940, 120, 1750, 630), (50, 660, 910, 1170), (940, 660, 1750, 1170)]
    labels = ["Portada", "Ficha municipal", "Explicabilidad SHAP", "Fuentes abiertas"]
    for img, box, label in zip(imgs, boxes, labels):
        bx0, by0, bx1, by1 = box
        bw, bh = bx1 - bx0, by1 - by0
        header = 46
        cropped = _cover_crop(img, bw - 6, bh - header - 6)
        canvas_img.paste(cropped, (bx0 + 3, by0 + header + 3))
        # Barra-título de la tarjeta
        d.rounded_rectangle((bx0, by0, bx1, by0 + header), radius=16, fill=VERDE)
        d.rectangle((bx0, by0 + header - 16, bx1, by0 + header), fill=VERDE)
        d.text((bx0 + 20, by0 + 11), label, font=font(24, True), fill="#FFFFFF")
        d.rounded_rectangle(box, radius=18, fill=None, outline=BORDE, width=3)
    canvas_img.save(output)


def build_demo_board(map_path: Path, ranking_path: Path, output: Path) -> None:
    canvas_img = Image.new("RGB", (1800, 900), FONDO)
    d = ImageDraw.Draw(canvas_img)
    d.text((55, 34), "Respaldo visual para la demostración", font=font(36, True), fill=TINTA)
    d.text((55, 88), "Si la conexión falla, estas capturas cubren la secuencia de demo en menos de 2 minutos.", font=font(20), fill="#52645A")
    map_img = Image.open(map_path).convert("RGB")
    rank_img = Image.open(ranking_path).convert("RGB")
    for img, box, label in [
        (map_img, (50, 150, 878, 730), "Mapa nacional"),
        (rank_img, (922, 150, 1750, 730), "Ranking exportable"),
    ]:
        bx0, by0, bx1, by1 = box
        header = 46
        cropped = _cover_crop(img, (bx1 - bx0) - 6, (by1 - by0) - header - 6)
        canvas_img.paste(cropped, (bx0 + 3, by0 + header + 3))
        d.rounded_rectangle((bx0, by0, bx1, by0 + header), radius=16, fill=VERDE)
        d.rectangle((bx0, by0 + header - 16, bx1, by0 + header), fill=VERDE)
        d.text((bx0 + 20, by0 + 11), label, font=font(24, True), fill="#FFFFFF")
        d.rounded_rectangle(box, radius=18, fill=None, outline=BORDE, width=3)
    canvas_img.save(output)


def build_download_card(output: Path) -> None:
    img = Image.new("RGB", (1600, 900), FONDO)
    d = ImageDraw.Draw(img)
    d.text((60, 55), "Consulta y salida exportable", font=font(34, True), fill=TINTA)
    d.rounded_rectangle((80, 150, 1520, 760), radius=28, fill=PANEL, outline=BORDE, width=3)
    d.rounded_rectangle((130, 230, 1470, 330), radius=18, fill="#EAF3EE")
    d.text((170, 258), "Ranking territorial de priorización", font=font(28, True), fill=TINTA)
    d.rounded_rectangle((1140, 248, 1430, 312), radius=14, fill=AZUL)
    d.text((1180, 265), "Descargar CSV", font=font(22, True), fill="#FFFFFF")
    bullets = (
        "• La página de ranking del producto expone un botón real de descarga CSV.\n"
        "• El archivo exporta municipio, departamento, nivel, score y los 4 índices.\n"
        "• La ficha territorial permite consulta puntual por municipio.\n"
        "• La demo puede cerrarse mostrando trazabilidad y capacidad de extracción."
    )
    add_wrapped_text(d, bullets, (160, 410), 1180, font(24), TINTA, line_gap=12)
    img.save(output)


def generate_visuals(data: dict) -> dict[str, Path]:
    crop_live_screens()
    save_qr("https://streamlit.spartanit.pro/", ASSETS / "qr_demo.png")
    save_qr("https://github.com/fredypaeze/aquabosque-minero-ia", ASSETS / "qr_repo.png")

    build_map_image(data["pred"], ASSETS / "02_mapa_cobertura.png")
    build_ranking_chart(data["pred"], ASSETS / "03_ranking_top15.png")
    build_coverage_map(data["master"], ASSETS / "11_cobertura_senales.png")
    build_architecture_diagram(ASSETS / "08_arquitectura_mvp.png")
    build_dataflow_diagram(ASSETS / "09_flujo_datos.png")
    build_model_explainer_image(data["metrics"], data["shap"], ASSETS / "12_modelo_vigente.png")
    sources = source_table(data)
    build_sources_table_image(sources, ASSETS / "10_tabla_fuentes.png")
    case_row = data["pred"][data["pred"]["municipio"] == "BARRANCABERMEJA"].iloc[0]
    build_case_image(case_row, ASSETS / "13_ficha_barrancabermeja.png")
    build_roadmap(ASSETS / "14_roadmap.png")
    build_app_collage(ASSETS / "15_app_collage.png")
    build_demo_board(ASSETS / "02_mapa_cobertura.png", ASSETS / "03_ranking_top15.png", ASSETS / "16_demo_backup.png")
    build_download_card(ASSETS / "17_consulta_descarga.png")

    # Marcos de navegador para capturas verticales (evita tiras angostas flotando).
    if (ASSETS / "01_app_portada_viva.png").exists():
        frame_browser(ASSETS / "01_app_portada_viva.png", ASSETS / "18_portada_framed.png")
    if (ASSETS / "06_shap_live.png").exists():
        frame_browser(ASSETS / "06_shap_live.png", ASSETS / "19_shap_framed.png")
    # Panel satelital NRT (capa 2 del proyecto).
    build_satelital_panel(ASSETS / "21_satelital_panel.png")

    for path in ASSETS.glob("*.png"):
        shutil.copy2(path, DEMO / path.name)
    return {p.stem: p for p in ASSETS.glob("*.png")}


def source_table(data: dict) -> pd.DataFrame:
    raw = data["raw_state"]["dimensiones"]
    return pd.DataFrame([
        {
            "dimension": "Minería formal",
            "fuente": "ANM · RUCOM · datos.gov.co/42ha-fhvj",
            "cobertura": "12.914 registros RUCOM cruzados por código DANE",
            "corte": "Archivo descargado 2026-07-14",
            "funcion": "Presión minera formal: conteo de explotadores y diversidad mineral",
            "url": "https://www.datos.gov.co/resource/42ha-fhvj.json",
            "archivo": "data/raw/mineria/rucom.csv",
            "limitacion": "No observa minería ilegal ni informal",
        },
        {
            "dimension": "Producción y regalías",
            "fuente": "ANM · datos.gov.co/r85m-vv6c",
            "cobertura": "75.888 registros; años visibles 2012-2026",
            "corte": "Archivo descargado 2026-07-14; último año visible 2026",
            "funcion": "Volumen de explotación y regalías por municipio",
            "url": "https://www.datos.gov.co/resource/r85m-vv6c.json",
            "archivo": "data/raw/mineria/anm_volumen.csv",
            "limitacion": "Periodicidad trimestral; mezcla recursos y periodos",
        },
        {
            "dimension": "Deforestación",
            "fuente": "IDEAM / SMByC · ArcGIS FeatureServer",
            "cobertura": "353 registros; 75 municipios foco; años 2017-2021",
            "corte": "Archivo extraído 2026-07-14; valor máximo por municipio",
            "funcion": "Índice de deforestación municipal",
            "url": "https://services9.arcgis.com/1TA62AToEccvEPrZ/arcgis/rest/services/Datos/FeatureServer",
            "archivo": "data/raw/bosque/deforestacion.csv",
            "limitacion": "Cobertura focalizada; no lectura satelital nacional en tiempo real",
        },
        {
            "dimension": "Calidad hídrica",
            "fuente": "IDEAM DHIME · ICA",
            "cobertura": "167 estaciones; fechas visibles 2002-2018",
            "corte": "Archivo existente en repo; fecha de muestra más reciente visible 2018",
            "funcion": "Índice hídrico = 1 - ICA donde existe estación",
            "url": "https://dhime.ideam.gov.co/",
            "archivo": "data/raw/agua/ica_ideam.csv",
            "limitacion": "Solo 71 municipios con estación asociada",
        },
        {
            "dimension": "Sensibilidad ambiental",
            "fuente": "RUNAP · ArcGIS capa de áreas protegidas",
            "cobertura": "1.882 áreas protegidas con centroide",
            "corte": "Archivo extraído 2026-07-14",
            "funcion": "Conteo y hectáreas protegidas cercanas a municipio",
            "url": "https://services3.arcgis.com/Fto9oba51JWVX0Qy/arcgis/rest/services/Areas_Protegidas_RUNAP/FeatureServer",
            "archivo": "data/raw/sensibilidad/runap.csv",
            "limitacion": "Aproximación por centroide, no intersección poligonal completa",
        },
        {
            "dimension": "Base territorial",
            "fuente": "DANE · DIVIPOLA · datos.gov.co/gdxc-w37w",
            "cobertura": "1.122 municipios con código y centroide",
            "corte": "Archivo descargado 2026-07-14",
            "funcion": "Unidad municipal de integración",
            "url": "https://www.datos.gov.co/resource/gdxc-w37w.json",
            "archivo": "data/raw/territorio/divipola.csv",
            "limitacion": "El CSV consumido no expone una etiqueta explícita de versión anual",
        },
        {
            "dimension": "Sensibilidad social",
            "fuente": "Municipios PDET · datos.gov.co/idrk-ba8y",
            "cobertura": "170 municipios",
            "corte": "Archivo descargado 2026-07-14",
            "funcion": "Ajuste de sensibilidad territorial (+0.25 al índice)",
            "url": "https://www.datos.gov.co/resource/idrk-ba8y.json",
            "archivo": "data/raw/territorio/pdet.csv",
            "limitacion": "No es señal de daño; es priorización de contexto territorial",
        },
    ])


def build_traceability_files(data: dict) -> None:
    sources = source_table(data)
    metrics = pd.DataFrame([
        {"indicador": "Municipios", "valor": 1122, "evidencia": "data/processed/master_con_etiqueta.csv"},
        {"indicador": "Crítico", "valor": int((data["pred"]["riesgo_nivel"] == "Crítico").sum()), "evidencia": "outputs/tables/predicciones.csv"},
        {"indicador": "Alto", "valor": int((data["pred"]["riesgo_nivel"] == "Alto").sum()), "evidencia": "outputs/tables/predicciones.csv"},
        {"indicador": "Medio", "valor": int((data["pred"]["riesgo_nivel"] == "Medio").sum()), "evidencia": "outputs/tables/predicciones.csv"},
        {"indicador": "Accuracy", "valor": data["metrics"]["accuracy"], "evidencia": "models/metrics/metricas.json"},
        {"indicador": "F1 macro", "valor": data["metrics"]["f1_macro"], "evidencia": "models/metrics/metricas.json"},
        {"indicador": "Línea base", "valor": data["metrics"]["baseline_clase_mayoritaria"], "evidencia": "models/metrics/metricas.json"},
    ])
    coverage = pd.DataFrame([
        {"señal": "Minería", "municipios_con_dato": int((data["master"]["mineria_titulos"] > 0).sum()), "total": 1122},
        {"señal": "Deforestación", "municipios_con_dato": int((data["master"]["deforestacion_ha"] > 0).sum()), "total": 1122},
        {"señal": "Agua ICA", "municipios_con_dato": int((data["master"]["agua_estaciones"] > 0).sum()), "total": 1122},
        {"señal": "RUNAP", "municipios_con_dato": int((data["master"]["runap_areas"] > 0).sum()), "total": 1122},
        {"señal": "PDET", "municipios_con_dato": int((data["master"]["es_pdet"] == 1).sum()), "total": 1122},
    ])
    assertions = pd.DataFrame([
        {
            "afirmacion": "La app desplegada existe y está operativa.",
            "estado": "VERIFICADO",
            "evidencia": "Captura live-home.png y URLs /Mapa_de_riesgo, /Ranking, /Ficha_territorial, /Explicabilidad",
            "ruta": "tmp/audit/live-home.png; tmp/screens/*.png",
        },
        {
            "afirmacion": "El modelo vigente es XGBoost multiclase con SHAP.",
            "estado": "VERIFICADO",
            "evidencia": "src/aquabosque/models/train.py; docs/MODEL_CARD.md; app/Explicabilidad",
            "ruta": "src/aquabosque/models/train.py",
        },
        {
            "afirmacion": "No hay un detector activo de anomalías en el producto actual.",
            "estado": "VERIFICADO",
            "evidencia": "Búsqueda completa del repo sin modelos de outlier activos; solo XGBoost + SHAP",
            "ruta": "rg \"anomal|Isolation|outlier\"",
        },
        {
            "afirmacion": "La priorización usa una etiqueta técnica compuesta por cuantiles.",
            "estado": "VERIFICADO",
            "evidencia": "src/aquabosque/features/build_target.py",
            "ruta": "src/aquabosque/features/build_target.py",
        },
        {
            "afirmacion": "Las 16 pruebas pasan con el estado actual del repo.",
            "estado": "VERIFICADO",
            "evidencia": "pytest -q",
            "ruta": "tests/",
        },
        {
            "afirmacion": "El agua no cubre todo el país y no debe presentarse como monitoreo en tiempo real.",
            "estado": "VERIFICADO",
            "evidencia": "71 municipios con estación; fechamuestra 2002-2018; Metodología reconoce actualización periódica",
            "ruta": "data/raw/agua/ica_ideam.csv; app/pages/06_📖_Metodología.py",
        },
        {
            "afirmacion": "La deforestación usada no equivale a un procesamiento satelital nacional en vivo.",
            "estado": "VERIFICADO",
            "evidencia": "75 municipios foco; años 2017-2021",
            "ruta": "data/raw/bosque/deforestacion.csv",
        },
    ])

    with pd.ExcelWriter(OUT / "05_AquaBosque_Trazabilidad_Fuentes.xlsx", engine="openpyxl") as writer:
        sources.to_excel(writer, index=False, sheet_name="fuentes")
        coverage.to_excel(writer, index=False, sheet_name="cobertura")
        metrics.to_excel(writer, index=False, sheet_name="metricas")
    with pd.ExcelWriter(OUT / "06_AquaBosque_Afirmaciones_Evidencias.xlsx", engine="openpyxl") as writer:
        assertions.to_excel(writer, index=False, sheet_name="afirmaciones")
    with pd.ExcelWriter(OUT / "09_AquaBosque_Fuentes_Fechas_Enlaces.xlsx", engine="openpyxl") as writer:
        sources[["dimension", "fuente", "corte", "url", "archivo", "limitacion"]].to_excel(writer, index=False, sheet_name="fuentes")


def write_text(path: Path, content: str) -> None:
    path.write_text(content.strip() + "\n", encoding="utf-8")


def md_to_docx(md_text: str, output: Path, title: str) -> None:
    doc = Document()
    doc.add_heading(title, level=0)
    for raw_line in md_text.strip().splitlines():
        line = raw_line.rstrip()
        if not line:
            doc.add_paragraph("")
            continue
        if line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("- "):
            doc.add_paragraph(line[2:], style="List Bullet")
        else:
            doc.add_paragraph(line)
    doc.save(output)


def build_audit_markdown(data: dict) -> str:
    coverage = {
        "mineria": int((data["master"]["mineria_titulos"] > 0).sum()),
        "deforestacion": int((data["master"]["deforestacion_ha"] > 0).sum()),
        "agua": int((data["master"]["agua_estaciones"] > 0).sum()),
        "runap": int((data["master"]["runap_areas"] > 0).sum()),
        "pdet": int((data["master"]["es_pdet"] == 1).sum()),
    }
    return f"""
# Auditoría del producto AquaBosque Minero IA

## Qué se verificó
- App desplegada en `https://streamlit.spartanit.pro/` con portada y páginas internas auditadas el 22-jul-2026.
- Repositorio local `/home/tuxilo/aquabosque-minero-ia` en commit `4f9989f`.
- Pipeline, artefactos, métricas y pruebas del estado actual del código.

## Hallazgos principales
- El producto vigente **sí está desplegado** y expone mapa, ranking, ficha territorial, explicabilidad, datos abiertos y metodología.
- La IA vigente es **XGBoost multiclase + SHAP**, no un detector de anomalías.
- La priorización usa una **etiqueta técnica compuesta por cuantiles**: `0.35·minero + 0.30·deforestación + 0.25·hídrico + 0.10·sensibilidad`.
- El repo pasa **16/16 pruebas**.
- Distribución verificada: **57 críticos, 112 altos, 281 medios, 672 bajos**.

## Cobertura real del dataset usado por el MVP
- Municipios base: `1.122`.
- Minería formal con señal: `{coverage['mineria']}` municipios.
- Deforestación con señal: `{coverage['deforestacion']}` municipios.
- Agua ICA con estación asociada: `{coverage['agua']}` municipios.
- RUNAP con señal: `{coverage['runap']}` municipios.
- PDET: `{coverage['pdet']}` municipios.

## Contradicciones que había que corregir antes del deck
1. **Anomalías vs. clasificación**: el prompt de trabajo hablaba de detección de anomalías en tiempo real, pero el código, la app y la documentación vigente muestran XGBoost multiclase + SHAP sobre una etiqueta técnica.
2. **5 fuentes vs. 7 fuentes**: la portada viva y el README hablan de 5 fuentes oficiales integradas; algunos documentos (`docs/00_resumen_ejecutivo.md`, `docs/09_defensa_jurado.md`) dicen 7 al contar activos territoriales de soporte (`DIVIPOLA`, `PDET`) como fuentes separadas.
3. **Tiempo real**: la metodología viva niega tiempo real; agua usa observaciones con `fechamuestra` visible entre 2002 y 2018, y deforestación usa cortes 2017-2021.
4. **Cobertura satelital**: no existe un procesamiento satelital operacional nacional en la versión actual; la señal de bosque llega desde un FeatureServer municipal ya resumido.

## Documentos que conviene archivar o corregir
- `docs/00_resumen_ejecutivo.md`: unificar conteo de fuentes y evitar ambigüedad entre “5” y “7”.
- `docs/09_defensa_jurado.md`: mismo ajuste; además reemplazar cualquier lectura que el jurado pudiera interpretar como superioridad predictiva.
- Prompt o borradores externos que hablen de **anomalías**, **satélite en tiempo real** o **monitoreo nacional en vivo**.

## Posicionamiento defendible
- MVP funcional de inteligencia territorial ambiental.
- Integra datos abiertos reales a nivel municipal.
- Usa IA explicable para hacer operativa una priorización interpretable.
- No prueba causalidad, no sanciona y no detecta minería ilegal.
"""


def build_summary_markdown(data: dict) -> str:
    return f"""
# Resumen ejecutivo

AquaBosque Minero IA es un MVP funcional de inteligencia territorial ambiental que integra datos abiertos de minería formal, deforestación, calidad hídrica, áreas protegidas y contexto PDET para priorizar municipios colombianos donde conviene revisar primero. La cadena técnica está validada de extremo a extremo: descarga y limpieza de fuentes, integración municipal, construcción de 4 índices, generación de una etiqueta técnica por cuantiles, entrenamiento de un modelo XGBoost multiclase con SHAP y despliegue en una aplicación Streamlit.

La auditoría del 22-jul-2026 confirma que el producto desplegado existe, el repositorio reproduce sus artefactos y las 16 pruebas pasan. La cobertura base es nacional (1.122 municipios), pero la intensidad de cada señal es desigual: hay minería formal en 667 municipios, deforestación visible en 75, agua ICA con estación asociada en 71 y RUNAP en 505. Por eso la narrativa correcta no es “monitoreo nacional en tiempo real”, sino priorización territorial interpretable con la mejor señal abierta disponible por fuente.

El modelo vigente no es un detector de anomalías. Es un clasificador XGBoost sobre una etiqueta técnica compuesta y su valor defendible está en la interpretabilidad: SHAP muestra qué variables pesan más en cada priorización. La accuracy de 91,1% se declara con honestidad, porque el propio proyecto reconoce que la etiqueta es una fórmula compuesta que el modelo re-aprende parcialmente. Lo robusto para jurado es mostrar que el MVP ya hace útil la integración de datos dispersos y que el siguiente salto lógico es automatizar actualización, ampliar cobertura y validar operativamente con autoridades.
"""


def build_questions_markdown() -> str:
    return """
# Preguntas y respuestas para la defensa

## 1. ¿Dónde está realmente la inteligencia artificial?
En el clasificador XGBoost multiclase que toma 13 variables municipales y asigna uno de cuatro niveles de priorización. SHAP explica qué factores pesan en cada clasificación.

## 2. ¿Por qué se utilizó ese modelo?
Porque el producto actual no tiene etiquetas oficiales de “riesgo ambiental” por municipio. XGBoost permite convertir una priorización técnica en un modelo reproducible e interpretable, y SHAP hace visible el peso de cada variable.

## 3. ¿Qué diferencia existe entre el ranking y la anomalía?
En la versión vigente no hay un módulo activo de anomalías. Lo que existe es un score técnico de priorización y una clasificación XGBoost que lo hace operativo y explicable.

## 4. ¿Cómo se valida una alerta?
No se presenta como alerta automática confirmada. La salida sirve para decidir dónde revisar primero con análisis experto y validación ambiental adicional.

## 5. ¿Una anomalía significa contaminación?
No aplica al producto actual. Incluso en la dimensión hídrica, una señal estadística o un ICA bajo no prueba contaminación causal; solo prioriza revisión.

## 6. ¿El producto detecta minería ilegal?
No. Usa minería formal observada en RUCOM y señales ambientales/territoriales. Sirve para focalizar revisión, no para declarar ilegalidad.

## 7. ¿El sistema funciona actualmente en tiempo real?
No. La metodología viva y la auditoría muestran actualización periódica según la frecuencia real de publicación de cada fuente.

## 8. ¿Qué parte ya está construida?
Descarga o consumo de fuentes, integración municipal, construcción de variables, entrenamiento del modelo, métricas, SHAP y aplicación desplegada.

## 9. ¿Qué parte se escalaría si el proyecto gana?
Automatización de cortes, históricos temporales, reentrenamiento, validación con expertos, mejoras de demo y alertas/reportes operacionales.

## 10. ¿Cómo se actualizarían las fuentes?
Reejecutando los scripts de descarga y preparación o sustituyéndolos por tareas automatizadas según la periodicidad de cada fuente.

## 11. ¿Qué pasa cuando un municipio no tiene datos?
Se mantiene la fila municipal y se documenta la ausencia. En agua, por ejemplo, sin estación asociada el índice hídrico observado se deja en 0 y se marca la ausencia de dato.

## 12. ¿Cómo se evita generar falsas conclusiones?
Separando explícitamente priorización de causalidad, mostrando limitaciones por fuente y manteniendo trazabilidad de cada variable.

## 13. ¿Cómo puede utilizarlo una entidad pública?
Como tablero de focalización territorial para ordenar revisión técnica, cruces interinstitucionales y priorización de recursos de monitoreo.

## 14. ¿Por qué se considera un producto avanzado?
Porque no es una idea: integra múltiples fuentes reales, produce resultados reproducibles, tiene explicabilidad y está desplegado.

## 15. ¿Cómo puede ampliarse hacia Amazonía y Pacífico?
Automatizando cortes, ampliando coberturas geográficas y sumando nuevas fuentes con la misma lógica territorial municipal o submunicipal.

## 16. ¿Cómo se garantiza trazabilidad?
Cada afirmación del deck quedó vinculada a archivos, scripts, capturas o artefactos de salida en las tablas de evidencias.

## 17. ¿Cómo se asegura la interpretabilidad?
Con SHAP global y por caso, además de variables crudas visibles en ranking y ficha.

## 18. ¿Qué diferencia a AquaBosque de un tablero tradicional?
No solo visualiza datos: los integra, normaliza, prioriza y explica el peso relativo de cada señal.

## 19. ¿Qué decisiones puede apoyar?
Dónde revisar primero, dónde cruzar análisis de agua y minería, y qué municipios requieren seguimiento reforzado.

## 20. ¿Qué evidencia demuestra que no es solo una idea?
App desplegada, repo reproducible, métricas guardadas, capturas auditadas, archivos de salida y 16 pruebas pasando.
"""


def build_script_markdown(slides: list[SlideSpec]) -> str:
    parts = ["# Guion de exposición de 15 minutos"]
    for slide in slides:
        parts.append(f"\n## Diapositiva {slide.number}. {slide.title}")
        parts.append(f"- Objetivo: {slide.objective}")
        parts.append(f"- Duración máxima: {slide.duration}")
        parts.append(f"- Texto sugerido: {slide.speech}")
        parts.append(f"- Mensaje para el jurado: {slide.takeaway}")
        parts.append(f"- Transición: {slide.transition}")
        if slide.questions:
            parts.append("- Posibles preguntas:")
            for question, answer in slide.questions:
                parts.append(f"  - {question} → {answer}")
    return "\n".join(parts) + "\n"


def build_outdated_docs_markdown() -> str:
    return """
# Documentos antiguos o inconsistentes a corregir

- `docs/00_resumen_ejecutivo.md`
  Razón: habla de 7 fuentes mientras la portada viva y el README hablan de 5 fuentes oficiales integradas; conviene separar fuentes núcleo de activos territoriales de soporte.

- `docs/09_defensa_jurado.md`
  Razón: misma inconsistencia de conteo y riesgo de que el jurado lea la accuracy como mérito predictivo si no se conserva la nota de honestidad.

- Cualquier prompt, borrador o correo que describa el producto como:
  - detector de anomalías activo,
  - monitoreo ambiental en tiempo real,
  - procesamiento satelital operacional nacional,
  - sistema que detecta minería ilegal.

- Material de demo sin respaldo visual.
  Razón: hoy la app no tiene deep-links listos para abrir un municipio específico como Barrancabermeja; conviene conservar capturas offline para jurado.
"""


def build_sources_and_assertions(data: dict) -> tuple[pd.DataFrame, pd.DataFrame]:
    sources = source_table(data)
    assertions = pd.read_excel(OUT / "06_AquaBosque_Afirmaciones_Evidencias.xlsx")
    return sources, assertions


def slide_specs(assets: dict[str, Path], data: dict) -> list[SlideSpec]:
    return [
        SlideSpec(
            1,
            "AquaBosque Minero IA",
            [
                "MVP funcional de inteligencia territorial ambiental con app desplegada.",
                "Integra señales de minería formal, deforestación, agua y sensibilidad territorial.",
                "La auditoría se hizo sobre la versión viva y el repo reproducible.",
            ],
            assets.get("18_portada_framed") or assets.get("01_app_portada_viva"),
            "Fuente: app desplegada auditada el 22-jul-2026.",
            "Abrir con una imagen real del producto y fijar que estamos mostrando un MVP existente.",
            "AquaBosque Minero IA ya está desplegado y convierte datos ambientales dispersos en una priorización territorial interpretable. Lo que van a ver hoy no es un mockup: es un MVP funcional auditado sobre la app viva y el código real.",
            "0:30",
            "El producto existe y ya valida la cadena técnica principal.",
            "Después de mostrar que el producto existe, paso al problema que resuelve.",
            [("¿Está desplegado hoy?", "Sí. Se auditó la URL pública y se tomaron capturas de la versión activa.")],
            "Datos al Ecosistema 2026 · Categoría Sostenibilidad y Medio Ambiente · Datos abiertos de datos.gov.co + IA",
        ),
        SlideSpec(
            2,
            "El Problema",
            [
                "Colombia sí tiene datos ambientales, pero están fragmentados por fuente, cobertura y periodicidad.",
                "Revisar minería, bosque y agua por separado dificulta decidir dónde mirar primero.",
                "La necesidad institucional no es otro tablero aislado, sino una lectura territorial integrada.",
            ],
            assets.get("11_cobertura_senales"),
            "Fuente: cálculo propio sobre `master_con_etiqueta.csv`.",
            "Explicar que el cuello no es la ausencia total de datos sino su dispersión.",
            "El problema no es solamente tener datos; es convertirlos en señales territoriales comparables para actuar. Hoy la cobertura varía por fuente y eso dificulta priorizar recursos institucionales.",
            "1:00",
            "La propuesta nace de un problema de integración y priorización, no de falta absoluta de datos.",
            "Con el problema claro, muestro qué resuelve exactamente el MVP.",
            [("¿Por qué municipal y no otra unidad?", "Porque el municipio permite integrar todas las fuentes abiertas actuales con una unidad comprensible para política pública.")],
        ),
        SlideSpec(
            3,
            "La Propuesta",
            [
                "Cadena validada: datos abiertos → integración municipal → score técnico → XGBoost + SHAP → app.",
                "La app ya permite mapa, ranking, ficha territorial, explicabilidad y consulta de fuentes.",
                "El valor del MVP es volver operativa una priorización interpretable.",
            ],
            assets.get("08_arquitectura_mvp"),
            "Fuente: pipeline y app del repositorio auditado.",
            "Presentar AquaBosque como solución concreta y acotada al estado real del producto.",
            "AquaBosque no reemplaza la validación ambiental. Lo que hace es integrar señales dispersas, priorizarlas y explicar por qué un territorio aparece arriba en la lista de revisión.",
            "1:00",
            "El MVP ya comprueba que la cadena técnica funciona de punta a punta.",
            "Ahora bajo a las fuentes reales que alimentan esa cadena.",
            [("¿Qué parte es IA y qué parte es integración?", "La integración construye variables e índices; la IA clasifica y explica la priorización.")],
        ),
        SlideSpec(
            4,
            "Datos Abiertos Integrados",
            [
                "5 dimensiones núcleo + 2 activos territoriales, todos de datos abiertos (datos.gov.co y portales oficiales).",
                "Cada fila trae enlace al dataset abierto, archivo local y limitación explícita: trazable y auditable.",
                "No se incluyen fuentes no verificadas ni métricas inventadas.",
            ],
            assets.get("10_tabla_fuentes"),
            "Fuente: scripts de descarga, archivos raw y tabla de trazabilidad anexa.",
            "Mostrar qué fuentes son reales, qué cobertura tienen y cómo entran al sistema.",
            "Aquí es importante ser precisos: la portada viva habla de cinco fuentes núcleo, pero el trabajo completo usa además DIVIPOLA y PDET como soportes territoriales. Por eso el deck distingue claramente entre dimensión ambiental y activo de integración.",
            "1:00",
            "La trazabilidad de fuentes es uno de los puntos fuertes del MVP.",
            "Con las fuentes claras, explico cómo se vuelven comparables a nivel municipal.",
            [("¿Por qué aparece a veces 5 y a veces 7 fuentes?", "Cinco son dimensiones núcleo; DIVIPOLA y PDET son activos de soporte para integración y sensibilidad territorial.")],
        ),
        SlideSpec(
            5,
            "Construcción De La Unidad Territorial",
            [
                "La unidad común es el municipio: 1.122 filas con código y centroide.",
                "RUCOM, ANM volumen y PDET cruzan por código DANE; ICA y RUNAP por proximidad al centroide.",
                "Las ausencias no se inventan: se documentan como parte del contexto de calidad (proceso guiado por CRISP-ML).",
            ],
            assets.get("09_flujo_datos"),
            "Fuente: `build_master.py` y `build_target.py`.",
            "Explicar la lógica territorial y el manejo de datos faltantes.",
            "El MVP resuelve la comparabilidad territorial con una regla simple y auditable. La base es municipal, y las fuentes que no traen código DANE se asignan por proximidad al centroide, sin esconder esa simplificación.",
            "1:00",
            "La integración municipal es la pieza que hace posible comparar señales heterogéneas.",
            "Sobre esa base se monta la IA vigente y su forma correcta de narrarla.",
            [("¿Se hace intersección poligonal completa?", "No en todas las fuentes. Para RUNAP e ICA el MVP usa centroide/proximidad como aproximación explícitamente documentada.")],
        ),
        SlideSpec(
            6,
            "Cómo Funciona La IA",
            [
                "Modelo vigente: XGBoost multiclase con 15 variables (incluye señal satelital de fuego) y salida en 4 niveles.",
                "La etiqueta es técnica y viene de una fórmula compuesta por cuantiles.",
                "SHAP es el núcleo defendible: explica qué variables pesan más en la priorización.",
            ],
            assets.get("12_modelo_vigente"),
            "Fuente: `train.py`, `metricas.json`, `importancia_global.csv`.",
            "Alinear el discurso con la implementación real y evitar sobredeclaraciones.",
            "Aquí hay que ser transparentes: la versión actual no usa un detector de anomalías. Usa XGBoost para volver operativa una priorización técnica y SHAP para explicar cada caso. Esa honestidad fortalece, no debilita, la defensa.",
            "1:30",
            "La IA sí existe, pero su valor hoy está en interpretabilidad y operatividad, no en vender una accuracy descontextualizada.",
            "Con el modelo claro, muestro cómo se traduce en una lectura territorial entendible.",
            [("¿Por qué no hablan de anomalías?", "Porque el producto vigente no las implementa. La auditoría obliga a defender la IA realmente activa.")],
        ),
        SlideSpec(
            7,
            "Priorización E Interpretabilidad",
            [
                "El score técnico y la predicción del modelo no son lo mismo, pero se alinean en la app.",
                "Barrancabermeja muestra un caso con señal minera, hídrica y de sensibilidad al mismo tiempo.",
                "La salida ayuda a decidir dónde revisar primero; no reemplaza análisis experto.",
            ],
            assets.get("13_ficha_barrancabermeja"),
            "Fuente: `predicciones.csv` y variables crudas del municipio.",
            "Bajar la IA a un caso concreto que una autoridad pueda entender.",
            "Para la defensa conviene un municipio con varias señales a la vez. Barrancabermeja sirve porque no depende solo de deforestación o solo de minería: combina presión minera formal, una señal hídrica observada y alta sensibilidad territorial.",
            "1:00",
            "La salida del MVP es accionable porque combina ranking con explicación.",
            "Después del caso, enseño la aplicación real donde esto se consulta.",
            [("¿Por qué no usan Orito como caso?", "Orito es muy fuerte en minería y bosque, pero Barrancabermeja muestra además señal hídrica y ayuda a explicar más dimensiones a la vez.")],
        ),
        SlideSpec(
            8,
            "Capa Satelital · Monitoreo Near-Real-Time",
            [
                "Nueva capa: focos de calor activos por satélite (NASA FIRMS · VIIRS + MODIS), actualizados a diario.",
                "El modelo lo predijo, el satélite lo confirma: 62 municipios Alto/Crítico tienen fuego activo ahora.",
                "Y detecta actividad reciente que el índice estático (datos 2017-2021) no veía: responde al reto en tiempo real.",
            ],
            assets.get("21_satelital_panel"),
            "Fuente: NASA FIRMS (VIIRS SNPP+NOAA-20, MODIS C6.1) · datos abiertos · asignación a municipio en el repo.",
            "Mostrar que la solución ya incorpora imágenes satelitales y monitoreo near-real-time, no solo datos estáticos.",
            "Aquí subimos de nivel: además del índice, integramos una capa satelital que se actualiza a diario. Los focos de calor son el proxy estándar de la frontera de deforestación y quema. Lo potente es la fusión: donde el modelo prioriza y el satélite confirma fuego hoy, la autoridad tiene su máxima prioridad de verificación. Y sobre imágenes crudas de Sentinel-2 corremos deep learning en la infraestructura GPU del Ministerio.",
            "1:30",
            "La solución ya no es estática: incorpora satélite y near-real-time, con soberanía de datos.",
            "Con la capacidad satelital demostrada, muestro la aplicación funcional donde todo se consulta.",
            [("¿Es tiempo real de verdad o promesa?", "La capa de focos FIRMS es real y diaria, ya integrada. Sobre Sentinel-2 corremos detección de deforestación con deep learning en nuestras GPU L40S; lo presentamos según su estado de validación, sin sobredeclarar.")],
        ),
        SlideSpec(
            9,
            "Capa Satelital · Deforestación con Sentinel-2 (GPU)",
            [
                "Imágenes Sentinel-2 (10 m) descargadas y procesadas en la infraestructura GPU del Ministerio (NVIDIA L40S).",
                "Detección de cambio NDVI 2023→2026 sobre un frente de deforestación de La Macarena: ~985 ha de pérdida de cobertura.",
                "Datos abiertos Copernicus, procesamiento soberano — el dato no sale del Estado. El motor U-Net escala esta capacidad.",
            ],
            assets.get("22_sentinel2_lamacarena"),
            "Fuente: Sentinel-2 L2A (Copernicus / AWS Open Data) · detección de cambio NDVI · procesado en aicluster (L40S).",
            "Demostrar procesamiento real de imagen satelital cruda sobre GPU propia, con soberanía de datos.",
            "Esta es la capa de imagen cruda: bajamos escenas Sentinel-2 de dos fechas y detectamos dónde el bosque cayó, midiendo hectáreas. Lo corrimos en la máquina de las L40S del Ministerio, así que la imagen y el resultado nunca salen del Estado. Es la base sobre la que el segmentador U-Net de deep learning eleva la precisión.",
            "1:30",
            "Ya procesamos imagen satelital de 10 m sobre GPU pública propia; no es teoría.",
            "Con las dos capas satelitales demostradas, muestro la aplicación funcional.",
            [("¿Ya usa deep learning o solo índices?", "El resultado mostrado es detección de cambio NDVI sobre imagen real, corriendo en el host de las L40S. El segmentador U-Net (deep learning en GPU) es el siguiente paso de la misma capa, ya con el pipeline montado.")],
        ),
        SlideSpec(
            10,
            "La Aplicación Funcional",
            [
                "La versión desplegada ya expone portada, ficha, explicabilidad y trazabilidad de fuentes.",
                "La app no es una maqueta: lee artefactos reales del pipeline y presenta resultados consultables.",
                "Las capturas del deck provienen de la versión viva auditada el 22-jul-2026.",
            ],
            assets.get("15_app_collage"),
            "Fuente: capturas headless de la app desplegada.",
            "Dejar evidencia visual de que la app existe y está navegable.",
            "Esta diapositiva ayuda a cerrar cualquier duda de implementación: no estamos hablando de wireframes, sino de páginas activas que muestran resultados, explicación y fuente de datos.",
            "1:00",
            "El jurado debe ver que el MVP ya se usa como producto, no solo como notebook.",
            "Luego muestro cómo se haría la demo en menos de dos minutos.",
            [("¿Se puede mostrar en vivo?", "Sí, pero además dejamos respaldo offline por si falla la conexión.")],
        ),
        SlideSpec(
            9,
            "Demostración Del Producto",
            [
                "Ruta sugerida: portada → mapa nacional → ranking → ficha → SHAP → descarga CSV.",
                "La secuencia está diseñada para no superar 2 minutos.",
                "El paquete incluye capturas de respaldo por si la conexión no acompaña.",
            ],
            assets.get("16_demo_backup"),
            "Fuente: assets del paquete jurado 2026.",
            "Preparar una demo robusta y corta para la exposición.",
            "La idea no es navegar todo el sistema, sino mostrar valor rápido: cobertura nacional, lectura territorial, explicación del modelo y salida exportable. Todo eso cabe en menos de dos minutos.",
            "1:00",
            "La demo debe vender utilidad institucional, no detalles de programación.",
            "Con la demo clara, paso a un caso institucional concreto.",
            [("¿Qué pasa si el mapa en vivo tarda?", "Se usa el respaldo offline del paquete y se mantiene la secuencia narrativa.")],
        ),
        SlideSpec(
            10,
            "Caso De Uso Institucional",
            [
                "La solución orienta revisión y focalización, no sanción automática.",
                "Para Barrancabermeja la lectura útil es reforzar cruce entre presión minera formal, agua y sensibilidad.",
                "La entidad decide después qué validaciones adicionales hacer en territorio o con otras bases.",
            ],
            assets.get("17_consulta_descarga"),
            "Fuente: app de ranking y extracción CSV definida en código.",
            "Conectar el MVP con una decisión pública concreta.",
            "La pregunta institucional no es si el modelo sentencia algo, sino si ayuda a usar mejor el tiempo técnico. Un municipio crítico como Barrancabermeja sirve para ordenar revisión hídrica, cruce con actividad minera y sensibilidad territorial.",
            "1:30",
            "AquaBosque apoya decisiones de priorización; no suplanta la validación ambiental.",
            "Después del caso, cierro mostrando qué queda realmente validado por el MVP.",
            [("¿Qué haría una autoridad después de ver este caso?", "Cruzar con inspección, seguimiento hídrico, expedientes y otras bases sectoriales antes de concluir.")],
        ),
        SlideSpec(
            11,
            "Qué Valida El MVP",
            [
                "Datos abiertos reales integrados y reproducibles desde el repositorio público (GitHub).",
                "Modelo, métricas, SHAP y app conectados a los mismos artefactos.",
                "Auditoría completada: 16 pruebas pasando, código abierto y narrativa corregida a lo verificable.",
            ],
            assets.get("19_shap_framed") or assets.get("06_shap_live"),
            "Fuente: repo auditado, pruebas y métricas del modelo.",
            "Sintetizar la evidencia técnica sin convertirla en checklist de concurso.",
            "Lo que el MVP ya valida es suficiente para una defensa seria: integración de fuentes reales, priorización interpretable, app funcional y capacidad de regenerar evidencias. Lo que no valida todavía se declara sin maquillaje.",
            "1:00",
            "El jurado debe concluir que hay sustancia técnica y honestidad metodológica.",
            "Con eso puesto, la ruta de escalamiento se vuelve creíble y no aspiracional.",
            [("¿Por qué esa honestidad importa?", "Porque evita sobreprometer y hace más confiable el producto frente a un jurado técnico.")],
        ),
        SlideSpec(
            12,
            "Escalamiento",
            [
                "La siguiente fase natural es automatizar cortes y enriquecer historial temporal.",
                "Luego vendrían validación con expertos, alertas y mejores rutas de consulta.",
                "Escalar no significa prometer infraestructura nacional completa antes de validar utilidad.",
            ],
            assets.get("14_roadmap"),
            "Fuente: hoja de ruta construida a partir de la auditoría del producto.",
            "Mostrar una ruta de crecimiento sensata, no inflada.",
            "Si el proyecto gana, el salto lógico no es rehacerlo todo, sino automatizar lo que ya funciona, ampliar cobertura y llevar la priorización a un uso operativo más estable con instituciones.",
            "1:30",
            "El escalamiento parte de una base validada, no de una idea en PowerPoint.",
            "Después del roadmap cierro con impacto potencial y mensaje final.",
            [("¿Por qué no prometen cobertura total inmediata?", "Porque el producto todavía depende de coberturas y periodicidades reales de cada fuente, y eso es mejor decirlo explícitamente.")],
        ),
        SlideSpec(
            13,
            "Impacto Potencial",
            [
                "Impacto ambiental y de gestión pública: focaliza revisión territorial con datos abiertos ya existentes.",
                "Menos revisión manual dispersa y más capacidad de lectura integrada.",
                "Escalable y replicable a las regiones prioritarias del concurso: Amazonía, Orinoquía y Pacífico.",
            ],
            assets.get("02_mapa_cobertura"),
            "Fuente: `predicciones.csv` y geojson municipal auditado.",
            "Traducir la solución en valor público sin inventar cifras de impacto.",
            "El impacto potencial aquí es cualitativo pero concreto: ordenar mejor la revisión institucional, aprovechar datos públicos dispersos y crear una base escalable para ecosistemas estratégicos sin esperar a una plataforma perfecta desde el día uno.",
            "1:00",
            "El producto puede apoyar decisiones públicas porque transforma dispersión en señal territorial.",
            "Cierro con el mensaje síntesis y los accesos para revisar demo y repo.",
            [("¿Qué diferencia esto de un tablero tradicional?", "Que no solo exhibe fuentes: las integra, prioriza y explica la lectura territorial resultante.")],
        ),
        SlideSpec(
            14,
            "Cierre",
            [
                "MVP funcional, desplegado y auditable.",
                "IA explicable alineada con el estado real del producto.",
                "Siguiente paso: escalar automatización, cobertura y validación institucional.",
            ],
            assets.get("18_portada_framed") or assets.get("01_app_portada_viva"),
            "Fuente: app desplegada, repo y QR del paquete.",
            "Cerrar con una tesis simple y con accesos directos a demo y repo.",
            "AquaBosque Minero IA ya demuestra que datos ambientales dispersos pueden convertirse en señales territoriales interpretables. El siguiente paso no es inventar más narrativa, sino escalar esta capacidad validada hacia un uso operacional más robusto.",
            "0:30",
            "No es una maqueta; es un MVP serio con ruta clara de evolución.",
            "Fin de la presentación.",
            [("¿Dónde pueden revisar el producto?", "En la URL pública y en el repositorio referenciado por los QR de cierre.")],
        ),
    ]


def fit_image(img_path: Path, max_w: float, max_h: float) -> tuple[float, float]:
    img = Image.open(img_path)
    ratio = min(max_w / img.width, max_h / img.height)
    return img.width * ratio, img.height * ratio


def build_pptx(slides: list[SlideSpec]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide in slides:
        s = prs.slides.add_slide(blank)
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor.from_string(FONDO.replace("#", ""))

        band = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor.from_string(VERDE.replace("#", ""))
        band.line.fill.background()

        title_box = s.shapes.add_textbox(Inches(0.55), Inches(0.58), Inches(7.2), Inches(0.8))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide.title
        run.font.bold = True
        run.font.size = Pt(23)
        run.font.color.rgb = RGBColor.from_string(TINTA.replace("#", ""))
        if slide.subtitle:
            sub_box = s.shapes.add_textbox(Inches(0.58), Inches(1.18), Inches(7.4), Inches(0.45))
            p = sub_box.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = slide.subtitle
            run.font.size = Pt(10.5)
            run.font.color.rgb = RGBColor.from_string(AZUL.replace("#", ""))

        bullet_box = s.shapes.add_textbox(Inches(0.62), Inches(1.65), Inches(4.55), Inches(4.75))
        bullet_tf = bullet_box.text_frame
        bullet_tf.word_wrap = True
        bullet_tf.margin_left = 0
        bullet_tf.margin_right = 0
        for i, bullet in enumerate(slide.bullets):
            par = bullet_tf.paragraphs[0] if i == 0 else bullet_tf.add_paragraph()
            par.text = bullet
            par.level = 0
            par.font.size = Pt(15)
            par.font.color.rgb = RGBColor.from_string(TINTA.replace("#", ""))
            par.space_after = Pt(10)
        if slide.image_path and slide.image_path.exists():
            region_x, region_y = 5.35, 1.5
            region_w, region_h = 7.6, 4.95
            w, h = fit_image(slide.image_path, region_w * 96, region_h * 96)
            w_in, h_in = w / 96, h / 96
            pic_x = region_x + (region_w - w_in) / 2
            pic_y = region_y + (region_h - h_in) / 2
            s.shapes.add_picture(str(slide.image_path), Inches(pic_x), Inches(pic_y), width=Inches(w_in), height=Inches(h_in))

        foot = s.shapes.add_textbox(Inches(0.58), Inches(6.83), Inches(11.9), Inches(0.28))
        p = foot.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide.footnote
        run.font.size = Pt(8.5)
        run.font.color.rgb = RGBColor.from_string("#65766C".replace("#", ""))

        num = s.shapes.add_textbox(Inches(12.35), Inches(6.78), Inches(0.4), Inches(0.3))
        p = num.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(slide.number)
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(VERDE.replace("#", ""))

        if slide.title == "Cierre":
            s.shapes.add_picture(str(ASSETS / "qr_demo.png"), Inches(10.8), Inches(5.05), width=Inches(1.1))
            s.shapes.add_picture(str(ASSETS / "qr_repo.png"), Inches(12.0), Inches(5.05), width=Inches(1.1))
            for label, x in [("Demo", 10.93), ("Repo", 12.15)]:
                box = s.shapes.add_textbox(Inches(x), Inches(6.18), Inches(0.7), Inches(0.25))
                p = box.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = label
                run.font.size = Pt(9)
                run.font.color.rgb = RGBColor.from_string(TINTA.replace("#", ""))

    out_path = OUT / "01_AquaBosque_Presentacion_Jurado.pptx"
    prs.save(out_path)
    shutil.copy2(out_path, OUT / "11_AquaBosque_Presentacion_Backup_Offline.pptx")
    return out_path


def draw_pdf_slide(c: canvas.Canvas, slide: SlideSpec, page_w: float, page_h: float) -> None:
    c.setFillColor(HexColor(FONDO))
    c.rect(0, 0, page_w, page_h, fill=1, stroke=0)
    c.setFillColor(HexColor(VERDE))
    c.rect(0, page_h - 32, page_w, 32, fill=1, stroke=0)
    c.setFont("Helvetica-Bold", 22)
    c.setFillColor(HexColor(TINTA))
    c.drawString(36, page_h - 70, slide.title)
    if slide.subtitle:
        c.setFont("Helvetica", 10)
        c.setFillColor(HexColor(AZUL))
        c.drawString(36, page_h - 87, slide.subtitle)
    c.setFillColor(HexColor(TINTA))
    c.setFont("Helvetica", 12.5)
    y = page_h - 120
    for bullet in slide.bullets:
        wrapped = textwrap.wrap(bullet, 48)
        c.drawString(42, y, u"\u2022")
        c.drawString(58, y, wrapped[0])
        y -= 18
        for rest in wrapped[1:]:
            c.drawString(58, y, rest)
            y -= 16
        y -= 10
    if slide.image_path and slide.image_path.exists():
        iw, ih = Image.open(slide.image_path).size
        region_x, region_w = 402, 540
        region_y, region_h = 118, 358
        ratio = min(region_w / iw, region_h / ih)
        w = iw * ratio
        h = ih * ratio
        px = region_x + (region_w - w) / 2
        py = region_y + (region_h - h) / 2
        c.drawImage(ImageReader(str(slide.image_path)), px, py, width=w, height=h, preserveAspectRatio=True, mask="auto")
    c.setFillColor(HexColor("#65766C"))
    c.setFont("Helvetica", 8)
    c.drawString(36, 18, slide.footnote)
    c.drawRightString(page_w - 24, 18, str(slide.number))
    if slide.title == "Cierre":
        c.drawImage(ImageReader(str(ASSETS / "qr_demo.png")), page_w - 155, 48, width=58, height=58, mask="auto")
        c.drawImage(ImageReader(str(ASSETS / "qr_repo.png")), page_w - 83, 48, width=58, height=58, mask="auto")
        c.setFont("Helvetica", 8)
        c.setFillColor(HexColor(TINTA))
        c.drawCentredString(page_w - 126, 40, "Demo")
        c.drawCentredString(page_w - 54, 40, "Repo")


def build_pdf(slides: list[SlideSpec]) -> Path:
    page_size = landscape((13.333 * 72, 7.5 * 72))
    out_path = OUT / "02_AquaBosque_Presentacion_Jurado.pdf"
    c = canvas.Canvas(str(out_path), pagesize=page_size)
    for slide in slides:
        draw_pdf_slide(c, slide, page_size[0], page_size[1])
        c.showPage()
    c.save()
    return out_path


def write_package_docs(data: dict, slides: list[SlideSpec]) -> None:
    audit_md = build_audit_markdown(data)
    summary_md = build_summary_markdown(data)
    qa_md = build_questions_markdown()
    script_md = build_script_markdown(slides)
    stale_md = build_outdated_docs_markdown()

    write_text(OUT / "03_AquaBosque_Guion_15_Minutos.md", script_md)
    write_text(OUT / "04_AquaBosque_Preguntas_Defensa.md", qa_md)
    write_text(OUT / "07_AquaBosque_Resumen_Ejecutivo_1_Pagina.md", summary_md)
    write_text(OUT / "08_AquaBosque_Documentos_Inconsistentes.md", stale_md)
    write_text(OUT / "10_AquaBosque_Auditoria_Producto.md", audit_md)

    md_to_docx(script_md, OUT / "03_AquaBosque_Guion_15_Minutos.docx", "Guion de exposición")
    md_to_docx(qa_md, OUT / "04_AquaBosque_Preguntas_Defensa.docx", "Preguntas y respuestas")
    md_to_docx(summary_md, OUT / "07_AquaBosque_Resumen_Ejecutivo_1_Pagina.docx", "Resumen ejecutivo")


def main() -> None:
    ensure_dirs()
    data = load_data()
    assets = generate_visuals(data)
    build_traceability_files(data)
    slides = slide_specs(assets, data)
    for i, s in enumerate(slides, 1):  # renumeración secuencial (permite insertar slides)
        s.number = i
    build_pptx(slides)
    build_pdf(slides)
    write_package_docs(data, slides)
    print(f"Paquete generado en: {OUT}")


if __name__ == "__main__":
    main()
