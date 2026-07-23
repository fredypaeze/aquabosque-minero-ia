"""Deck PREMIUM AquaBosque — PPTX nativo 100% editable, diseño editorial ejecutivo.

Textos de alto impacto en las slides; el guión completo en Notas del orador.
Tipografía Montserrat (títulos) + Open Sans (cuerpo); paleta esmeralda + azul
marino; imágenes con esquinas redondeadas y sombra suave; grid consistente.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "outputs" / "jurado_2026" / "assets"
OUT = ROOT / "outputs" / "jurado_2026" / "AquaBosque_Presentacion_PREMIUM.pptx"

# ---- Paleta ----
BG = RGBColor(0xF8, 0xF9, 0xFA)
NAVY = RGBColor(0x12, 0x3A, 0x5C)
EMER = RGBColor(0x1B, 0x7A, 0x46)
FOREST = RGBColor(0x14, 0x53, 0x2D)
INK = RGBColor(0x22, 0x30, 0x2A)
GRAY = RGBColor(0x5B, 0x6B, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINEA = RGBColor(0xDD, 0xE6, 0xE0)
CRIT = RGBColor(0xB9, 0x1C, 0x1C)
ALTO = RGBColor(0xEA, 0x58, 0x0C)
MEDIO = RGBColor(0xF5, 0x9E, 0x0B)
BAJO = RGBColor(0x16, 0xA3, 0x4A)
SOFT = RGBColor(0xEA, 0xF3, 0xEE)

F_TIT = "Montserrat"
F_BODY = "Open Sans"

SW, SH = Inches(13.333), Inches(7.5)


# ---------------- helpers de bajo nivel ----------------
def _run(p, text, size, color, font=F_BODY, bold=False, italic=False, spacing=None):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    if spacing is not None:
        rPr = r._r.get_or_add_rPr(); rPr.set("spc", str(spacing))
    return r


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tb, tf


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp._element.spPr)
    return sp


def _soft_shadow(spPr):
    spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}><a:outerShdw blurRad="90000" dist="38100" '
        f'dir="5400000" rotWithShape="0"><a:srgbClr val="0B1F17"><a:alpha val="20000"/>'
        f'</a:srgbClr></a:outerShdw></a:effectLst>'))


def image_card(slide, name, x, y, max_w, max_h, radius=4200):
    """Imagen con esquinas redondeadas + sombra suave, ajustada a (max_w,max_h)."""
    p = ASSETS / f"{name}.png"
    iw, ih = Image.open(p).size
    ratio = min(max_w / iw, max_h / ih)
    w, h = iw * ratio, ih * ratio
    px = x + (max_w - w) / 2
    pic = slide.shapes.add_picture(str(p), Inches(px), Inches(y), Inches(w), Inches(h))
    spPr = pic._element.spPr
    geom = spPr.find(qn("a:prstGeom"))
    if geom is not None:
        geom.set("prst", "roundRect")
        av = geom.find(qn("a:avLst"))
        av.append(parse_xml(f'<a:gd {nsdecls("a")} name="adj" fmla="val {radius}"/>'))
    _soft_shadow(spPr)
    return pic, w, h


def title_block(slide, kicker, titulo, ty=0.62):
    tb, tf = textbox(slide, 0.75, 0.42, 8.6, 0.4)
    _run(tf.paragraphs[0], kicker.upper(), 12.5, EMER, font=F_TIT, bold=True, spacing=180)
    tb, tf = textbox(slide, 0.72, ty + 0.26, 8.8, 1.6)
    for i, ln in enumerate(titulo.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = 1.04
        _run(p, ln, 33, NAVY, font=F_TIT, bold=True)


def bullets(slide, x, y, w, items, size=18.5, gap=14):
    tb, tf = textbox(slide, x, y, w, 4.6)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, "●  ", 11, EMER, font=F_BODY, bold=True)
        _run(p, it, size, INK, font=F_BODY)
        p.space_after = Pt(gap); p.line_spacing = 1.12
    return tb


def kpi(slide, x, y, numero, etiqueta, color=EMER, w=3.0):
    tb, tf = textbox(slide, x, y, w, 1.35)
    _run(tf.paragraphs[0], numero, 44, color, font=F_TIT, bold=True)
    p2 = tf.add_paragraph(); _run(p2, etiqueta, 12.5, GRAY, font=F_BODY); p2.space_before = Pt(2)
    return tb


def chip(slide, x, y, w, h, text, fill, tcolor=WHITE, size=13):
    sp = rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2); tf.margin_left = Pt(4); tf.margin_right = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _run(p, ln, size, tcolor, font=F_TIT, bold=True)
    return sp


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def base_slide(prs, footer=True, page=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, 13.333, 7.5, fill=BG)              # fondo
    rect(s, 0, 0, 13.333, 0.16, fill=EMER)            # barra superior acento
    if footer:
        tb, tf = textbox(s, 0.75, 7.06, 8, 0.3)
        _run(tf.paragraphs[0], "Ministerio de Minas y Energía", 10, GRAY, font=F_BODY)
    if page:
        tb, tf = textbox(s, 12.4, 7.04, 0.7, 0.3)
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        _run(tf.paragraphs[0], str(page), 10, GRAY, font=F_TIT, bold=True)
    return s


def kpi_row(slide, y, datos, x0=0.8, dx=2.75):
    for i, (num, lab, col) in enumerate(datos):
        kpi(slide, x0 + i * dx, y, num, lab, color=col, w=dx - 0.15)


# ---------------- construcción ----------------
def build():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    IMG_X, IMG_W = 6.55, 6.1   # columna derecha de imagen

    # 1 · PORTADA
    s = base_slide(prs, footer=False, page=1)
    rect(s, 0, 0, 13.333, 0.16, fill=EMER)
    tb, tf = textbox(s, 0.75, 0.7, 8, 0.4)
    _run(tf.paragraphs[0], "DATOS AL ECOSISTEMA 2026 · SOSTENIBILIDAD Y MEDIO AMBIENTE", 12, EMER, font=F_TIT, bold=True, spacing=160)
    tb, tf = textbox(s, 0.72, 1.35, 5.7, 1.6)
    _run(tf.paragraphs[0], "AquaBosque", 46, NAVY, font=F_TIT, bold=True)
    p = tf.add_paragraph(); _run(p, "Minero IA", 46, EMER, font=F_TIT, bold=True)
    tb, tf = textbox(s, 0.75, 3.5, 5.5, 1.8)
    _run(tf.paragraphs[0], "Inteligencia artificial que prioriza dónde actuar en Colombia por riesgo ambiental —minería, deforestación y agua— usando solo datos oficiales.", 18.5, INK, font=F_BODY)
    kpi_row(s, 5.35, [("1.122", "Municipios", NAVY), ("5", "Dimensiones", EMER), ("100%", "Datos abiertos", NAVY)], x0=0.8, dx=1.9)
    image_card(s, "18_portada_framed", 6.5, 1.35, 6.2, 4.9)
    tb, tf = textbox(s, 0.75, 7.06, 9, 0.3)
    _run(tf.paragraphs[0], "Ministerio de Minas y Energía · Grupo de Datos Estratégicos", 10, GRAY, font=F_BODY)
    notes(s, "Bienvenidos. Hoy presentamos AquaBosque Minero IA. En pocas palabras, es una herramienta inteligente que nos dice qué municipios de Colombia necesitan atención urgente por riesgos ambientales. Lo hace cruzando información de minería, deforestación y calidad del agua en una sola plataforma, usando únicamente datos oficiales.")

    # 2 · EL PROBLEMA
    s = base_slide(prs, page=2)
    title_block(s, "El problema", "No falta información.\nFalta integrarla.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Colombia tiene muchos datos ambientales, pero dispersos por fuente y por sector.",
        "Revisar minería, bosque y agua por separado retrasa las decisiones.",
        "El Ministerio necesita una lectura territorial integrada, no otro tablero aislado.",
    ])
    image_card(s, "11_cobertura_senales", IMG_X, 1.55, IMG_W, 4.9)
    notes(s, "El problema hoy no es que nos falte información. Colombia tiene muchísimos datos ambientales, pero están dispersos. Revisar por un lado la minería, por otro los bosques y por otro el agua, hace muy difícil saber dónde actuar primero. No venimos a proponer otro tablero aislado más, sino una lectura integrada del territorio para que el Ministerio pueda tomar decisiones rápidas.")

    # 3 · LA PROPUESTA
    s = base_slide(prs, page=3)
    title_block(s, "La propuesta", "De datos dispersos\na decisiones, hoy.")
    bullets(s, 0.78, 2.7, 5.5, [
        "Es un producto funcional, no un concepto.",
        "Une los datos a nivel municipal y la IA asigna un nivel de riesgo.",
        "Priorización operativa, transparente y fácil de entender.",
    ])
    pasos = [("Datos\nabiertos", EMER), ("Integración\nmunicipal", NAVY), ("IA +\nexplicación", EMER), ("App\nweb", NAVY)]
    px, pw, py = 6.6, 1.4, 2.9
    for i, (t, c) in enumerate(pasos):
        chip(s, px + i * 1.55, py, pw, 1.1, t, c, size=13.5)
        if i < 3:
            rect(s, px + i * 1.55 + pw + 0.02, py + 0.45, 0.13, 0.2, fill=LINEA, shape=MSO_SHAPE.CHEVRON)
    tb, tf = textbox(s, 6.6, 4.4, 6.0, 1.2)
    _run(tf.paragraphs[0], "Lo más valioso: vuelve OPERATIVA la priorización del territorio —transparente y entendible para cualquier analista o directivo.", 15, GRAY, font=F_BODY, italic=True)
    notes(s, "Nuestra solución ya es un producto funcional, no un concepto. Tomamos esos datos dispersos, los unimos a nivel municipal y una Inteligencia Artificial les asigna un nivel de riesgo. Lo más valioso de esta herramienta es que hace que la priorización de los territorios sea operativa, transparente y muy fácil de entender para cualquier analista o directivo.")

    # 4 · DATOS ABIERTOS
    s = base_slide(prs, page=4)
    title_block(s, "Fuentes", "Solo datos oficiales\ny verificables.")
    bullets(s, 0.78, 2.7, 5.3, [
        "Agencia Nacional de Minería, IDEAM, DANE y áreas protegidas.",
        "Cada dato es auditable y rastreable hasta su origen.",
        "Sin estimaciones sin fundamento ni cifras inventadas.",
    ])
    filas = [("Minería (títulos y volumen)", "ANM · datos.gov.co"),
             ("Deforestación", "IDEAM / SMByC"),
             ("Calidad del agua (ICA)", "IDEAM DHIME"),
             ("Áreas protegidas", "RUNAP"),
             ("Base municipal", "DANE DIVIPOLA")]
    ty = 1.65; rowh = 0.62
    chip(s, IMG_X, ty, IMG_W, 0.5, "7 FUENTES OFICIALES INTEGRADAS", NAVY, size=13)
    for i, (dim, ent) in enumerate(filas):
        yy = ty + 0.62 + i * rowh
        rect(s, IMG_X, yy, IMG_W, rowh - 0.1, fill=WHITE if i % 2 == 0 else SOFT, line=LINEA, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb, tf = textbox(s, IMG_X + 0.25, yy + 0.05, 3.7, rowh - 0.15, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], dim, 14, INK, font=F_BODY, bold=True)
        tb, tf = textbox(s, IMG_X + 4.0, yy + 0.05, 2.0, rowh - 0.15, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], ent, 12.5, EMER, font=F_BODY)
    notes(s, "Para que esto tenga validez institucional, la regla de oro fue usar solo fuentes oficiales y verificables. Integramos bases de la Agencia Nacional de Minería, el IDEAM y registros de áreas protegidas. Aquí no hay estimaciones sin fundamento ni datos inventados; todo es auditable y se puede rastrear hasta su origen.")

    # 5 · UNIDAD TERRITORIAL
    s = base_slide(prs, page=5)
    title_block(s, "Método", "Comparar peras con peras:\nel municipio.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Todo estandarizado a los 1.122 municipios del país.",
        "Los datos se cruzan por las coordenadas de cada lugar.",
        "Si falta un dato, se marca como ausencia. No se inventa.",
    ])
    kpi(s, 0.8, 5.35, "1.122", "municipios como unidad común", color=NAVY, w=5)
    pasos = [("1", "Base municipal\nDIVIPOLA"), ("2", "Cruces por\ncoordenada"), ("3", "Índices 0–1\npor dimensión"), ("4", "Priorización\nusable")]
    px, py = 6.6, 2.7
    for i, (n, t) in enumerate(pasos):
        cy = py + i * 1.02
        rect(s, px, cy, 0.62, 0.62, fill=EMER if i % 2 == 0 else NAVY, shape=MSO_SHAPE.OVAL)
        tb, tf = textbox(s, px, cy + 0.06, 0.62, 0.5, anchor=MSO_ANCHOR.MIDDLE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER; _run(tf.paragraphs[0], n, 20, WHITE, font=F_TIT, bold=True)
        tb, tf = textbox(s, px + 0.85, cy, 5.0, 0.9, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], t.replace("\n", "  "), 16, INK, font=F_BODY, bold=True)
    notes(s, "¿Cómo mezclamos peras con manzanas? Estandarizamos todo al nivel de los 1.122 municipios del país. Cruzamos los datos de minería y agua usando las coordenadas de cada lugar. Si en algún municipio no hay datos de calidad de agua, el sistema es honesto y lo marca como ausencia de dato, no inventa problemas donde no hay evidencia.")

    # 6 · CÓMO FUNCIONA LA IA
    s = base_slide(prs, page=6)
    title_block(s, "El modelo", "Clasifica el riesgo\ny explica el porqué.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Cuatro niveles: Bajo, Medio, Alto y Crítico.",
        "Lo valioso: explica qué factor pesó más (fuego, minería, deforestación).",
        "Inteligencia artificial interpretable, no una caja negra.",
    ])
    niveles = [("Bajo", BAJO), ("Medio", MEDIO), ("Alto", ALTO), ("Crítico", CRIT)]
    for i, (t, c) in enumerate(niveles):
        chip(s, 6.6 + i * 1.5, 2.7, 1.35, 0.7, t, c, size=14)
    image_card(s, "19_shap_framed", IMG_X, 3.65, IMG_W, 2.75)
    notes(s, "El motor detrás de esto clasifica el riesgo en cuatro niveles: Bajo, Medio, Alto y Crítico. Pero la verdadera magia de este modelo no es solo decirnos el nivel de riesgo, sino explicarnos el por qué. El sistema nos muestra exactamente qué factores —si fue el fuego, la minería o la deforestación— pesaron más para encender la alerta en un municipio específico.")

    # 7 · BARRANCABERMEJA
    s = base_slide(prs, page=7)
    title_block(s, "Caso real", "Barrancabermeja:\nriesgo Crítico, con evidencia.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Coinciden minería, afectación hídrica y sensibilidad ambiental.",
        "La IA no sanciona: dice “revisen aquí primero”.",
        "Evidencia combinada, lista para la autoridad.",
    ])
    chip(s, IMG_X, 1.7, 2.95, 0.95, "NIVEL\nCRÍTICO", CRIT, size=17)
    kpi(s, IMG_X + 3.2, 1.72, "0.43", "Score de priorización (0–1)", color=NAVY, w=2.9)
    dims = [("Minería", 0.73, FOREST), ("Hídrico", 0.43, NAVY), ("Sensibilidad", 0.71, EMER)]
    by = 3.05
    for i, (lab, val, c) in enumerate(dims):
        yy = by + i * 0.78
        tb, tf = textbox(s, IMG_X, yy, 2.0, 0.4); _run(tf.paragraphs[0], lab, 14, INK, font=F_BODY, bold=True)
        rect(s, IMG_X + 1.9, yy + 0.02, 3.6, 0.3, fill=SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, IMG_X + 1.9, yy + 0.02, 3.6 * val, 0.3, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb, tf = textbox(s, IMG_X + 5.55, yy, 0.55, 0.4); _run(tf.paragraphs[0], f"{val:.2f}", 13, GRAY, font=F_TIT, bold=True)
    notes(s, "Veamos un caso real: Barrancabermeja. El sistema lo marca en riesgo 'Crítico'. Al mirar el detalle, vemos que tiene actividad minera, afectación hídrica y alta sensibilidad ambiental al mismo tiempo. La IA no reemplaza al experto ni sanciona automáticamente; lo que hace es decirle a la autoridad: 'Revisen aquí primero, esta es la evidencia combinada'.")

    # 8 · SATÉLITE NRT
    s = base_slide(prs, page=8)
    title_block(s, "Capa satelital · tiempo real", "El modelo lo predice,\nel satélite lo confirma.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Focos de calor de la NASA, actualizados a diario.",
        "Riesgo alto + fuego hoy = prioridad máxima de verificación.",
        "Del análisis estático a la respuesta en tiempo real.",
    ])
    kpi_row(s, 5.15, [("296", "Municipios con fuego", ALTO), ("62", "Prioridad máxima", CRIT)], x0=0.8, dx=2.7)
    image_card(s, "20_mapa_focos_nrt", 7.4, 1.5, 4.6, 5.0)
    notes(s, "No nos quedamos solo con datos históricos. Le conectamos una señal de satélite de la NASA que detecta focos de calor casi en tiempo real. Si el sistema dice que un municipio es de alto riesgo, y hoy mismo el satélite detecta fuego allí, cruzamos esa información al instante. Es pasar del análisis estático a la respuesta en tiempo real.")

    # 9 · SENTINEL-2 GPU
    s = base_slide(prs, page=9)
    title_block(s, "Capa satelital · deep learning", "Zoom satelital: cuántas\nhectáreas se perdieron.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Imágenes Sentinel-2 procesadas en la GPU del Ministerio.",
        "Vemos el antes y el después: hectáreas exactas de pérdida.",
        "Soberanía total: el dato no sale del Estado.",
    ])
    kpi_row(s, 5.15, [("985", "Hectáreas detectadas", ALTO), ("0.77", "Precisión del modelo (IoU)", NAVY)], x0=0.8, dx=2.7)
    image_card(s, "22_sentinel2_lamacarena", IMG_X, 2.6, IMG_W, 3.3)
    notes(s, "También tenemos la capacidad de hacer un zoom satelital profundo usando imágenes procesadas en la infraestructura del Ministerio. Esto nos permite ver el antes y el después de una zona y calcular exactamente cuántas hectáreas de bosque se acaban de perder. Todo el procesamiento de las imágenes se hace en casa, garantizando la soberanía de nuestros datos.")

    # 10 · ASISTENTE
    s = base_slide(prs, page=10)
    title_block(s, "IA generativa soberana", "Pregúntale al sistema\nen lenguaje normal.")
    bullets(s, 0.78, 2.7, 5.4, [
        "“¿Por qué La Macarena está priorizada?” → respuesta con evidencia.",
        "No inventa: solo responde con lo que ya está recolectado.",
        "LLM local en las L40S; el dato no sale del país.",
    ])
    rect(s, IMG_X, 1.7, IMG_W, 0.85, fill=RGBColor(0xE8, 0xF0, 0xFE), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, IMG_X + 0.25, 1.82, IMG_W - 0.5, 0.6, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], "👤  ¿Por qué La Macarena aparece priorizada?", 15, NAVY, font=F_BODY, bold=True)
    rect(s, IMG_X, 2.7, IMG_W, 1.9, fill=WHITE, line=LINEA, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, IMG_X + 0.25, 2.85, IMG_W - 0.5, 1.65)
    _run(tf.paragraphs[0], "🤖  ", 13, EMER, font=F_BODY, bold=True)
    _run(tf.paragraphs[0], "Nivel Crítico (score 0.51): pesan una deforestación muy alta (0.97), fuego satelital (0.69) y sensibilidad máxima. Orienta revisión; no prueba causalidad.", 15.5, INK, font=F_BODY)
    chip(s, IMG_X, 4.8, IMG_W, 0.8, "Corre en las NVIDIA L40S del Ministerio · el dato no sale del Estado", EMER, size=13)
    notes(s, "Para facilitar aún más las cosas, integramos un asistente virtual. Usted puede preguntarle en lenguaje natural, por ejemplo: '¿Por qué La Macarena está priorizada?'. El asistente lee los datos del sistema y le explica la situación. Está configurado para no inventar nada; solo responde basado en la evidencia que ya tenemos recolectada.")

    # 11 · APLICACIÓN
    s = base_slide(prs, page=11)
    title_block(s, "El producto", "Una aplicación real,\nfácil de usar.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Mapa nacional de calor, ranking Top 15 y fichas por municipio.",
        "En menos de dos minutos encuentras lo que necesitas.",
        "Desplegada y navegable hoy: no es una maqueta.",
    ])
    image_card(s, "15_app_collage", IMG_X, 1.55, IMG_W, 4.9)
    notes(s, "Todo esto vive en una aplicación real y fácil de usar. Tenemos un mapa nacional de calor, un ranking con el 'Top 15' de municipios críticos y fichas detalladas por territorio. La navegación está diseñada para que en menos de dos minutos usted encuentre la información que necesita para justificar una acción.")

    # 12 · DEMOSTRACIÓN
    s = base_slide(prs, page=12)
    title_block(s, "Demostración", "Recorrido de valor\nen dos minutos.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Mapa → Ranking → Ficha → Explicación → Descarga.",
        "Secuencia pensada para mostrar valor rápido.",
        "Con respaldo offline por si falla la conexión.",
    ])
    image_card(s, "16_demo_backup", IMG_X, 2.15, IMG_W, 3.7)
    notes(s, "Todo esto vive en una aplicación real y fácil de usar. Tenemos un mapa nacional de calor, un ranking con el 'Top 15' de municipios críticos y fichas detalladas por territorio. La navegación está diseñada para que en menos de dos minutos usted encuentre la información que necesita para justificar una acción.")

    # 13 · DESCARGA
    s = base_slide(prs, page=13)
    title_block(s, "Uso institucional", "Los datos deben fluir:\nun clic a Excel o CSV.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Cualquier consulta o ranking se descarga al instante.",
        "Los equipos técnicos lo cruzan con sus reportes internos.",
        "Orienta la revisión; la decisión final es de la entidad.",
    ])
    rect(s, IMG_X, 2.0, IMG_W, 3.4, fill=WHITE, line=LINEA, line_w=1, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
    rect(s, IMG_X + 0.4, 2.45, IMG_W - 0.8, 0.75, fill=SOFT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(s, IMG_X + 0.65, 2.6, 3.5, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    _run(tf.paragraphs[0], "Ranking territorial de priorización", 15, INK, font=F_BODY, bold=True)
    chip(s, IMG_X + IMG_W - 2.05, 2.6, 1.6, 0.5, "⬇  Descargar CSV", NAVY, size=13)
    tb, tf = textbox(s, IMG_X + 0.45, 3.55, IMG_W - 0.9, 1.7)
    for i, t in enumerate(["Exporta municipio, departamento, nivel, score e índices.",
                            "Consulta puntual por municipio en la ficha.",
                            "Trazabilidad y capacidad de extracción para el equipo técnico."]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _run(p, "•  ", 13, EMER, bold=True); _run(p, t, 14.5, INK, font=F_BODY); p.space_after = Pt(8)
    notes(s, "Sabemos que los datos deben fluir. Por eso, cualquier consulta o ranking se puede descargar directamente en un archivo de Excel o CSV con un solo clic. El objetivo es que los equipos técnicos puedan llevarse esta información y cruzarla rápidamente con sus propios reportes internos.")

    # 14 · QUÉ VALIDA
    s = base_slide(prs, page=14)
    title_block(s, "Estado", "Auditado y funcionando.")
    checks = ["Datos reales integrados y reproducibles.",
              "Modelo auditable + código abierto en repositorio público.",
              "Satélite operando y aplicación desplegada.",
              "16 pruebas automáticas en verde."]
    tf = None
    ty = 2.7
    for i, t in enumerate(checks):
        yy = ty + i * 0.78
        rect(s, 0.8, yy, 0.5, 0.5, fill=EMER, shape=MSO_SHAPE.OVAL)
        tb, tf = textbox(s, 0.9, yy + 0.02, 0.4, 0.45, anchor=MSO_ANCHOR.MIDDLE)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER; _run(tf.paragraphs[0], "✓", 17, WHITE, font=F_TIT, bold=True)
        tb, tf = textbox(s, 1.5, yy, 10.5, 0.6, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], t, 20, INK, font=F_BODY)
    notes(s, "Esta versión ya fue auditada técnicamente y funciona. Nuestro siguiente paso lógico es automatizar que los datos se actualicen solos cada cierto tiempo y generar un sistema de alertas directas. No prometemos una infraestructura gigantesca desde el día uno, sino un crecimiento comprobable y útil.")

    # 15 · ESCALAMIENTO
    s = base_slide(prs, page=15)
    title_block(s, "Ruta", "Crecimiento comprobable,\nsin humo.")
    fases = [("Hoy · nivel avanzado", ["IA + satélite en GPU propia", "Incertidumbre calibrada (conformal)", "Detección de anomalías", "App desplegada"], EMER),
             ("Frontera técnica", ["Modelo espacial bayesiano (INLA/BYM)", "Deforestación RADD · Sentinel-1", "Actualización automática"], NAVY),
             ("Uso operacional", ["Alertas directas a la autoridad", "Inferencia causal (minería→agua)", "Validación experta"], FOREST)]
    fx, fw, fy = 0.8, 3.85, 2.7
    for i, (tit, items, c) in enumerate(fases):
        x = fx + i * (fw + 0.15)
        rect(s, x, fy, fw, 3.3, fill=WHITE, line=c, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=True)
        rect(s, x, fy, fw, 0.75, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb, tf = textbox(s, x + 0.25, fy + 0.12, fw - 0.5, 0.55, anchor=MSO_ANCHOR.MIDDLE)
        _run(tf.paragraphs[0], tit, 17, WHITE, font=F_TIT, bold=True)
        tb, tf = textbox(s, x + 0.3, fy + 1.0, fw - 0.6, 2.1)
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            _run(p, "•  ", 13, c, bold=True); _run(p, it, 15, INK, font=F_BODY); p.space_after = Pt(10)
    notes(s, "Y no nos quedamos en lo básico. Ya sumamos dos capacidades que muy pocas entidades en el país usan: "
             "incertidumbre calibrada con conformal prediction —decimos el nivel de un municipio con una garantía "
             "estadística del 90%, no un porcentaje suelto— y detección de anomalías, que de forma independiente "
             "confirma la priorización y descubre municipios atípicos que la fórmula no destaca. La frontera hacia "
             "la que avanzamos es clara y la conocemos: modelos espaciales bayesianos, deforestación por radar tipo "
             "Brasil, e inferencia causal. Crecimiento comprobable, sobre una base ya validada.")

    # 16 · IMPACTO
    s = base_slide(prs, page=16)
    title_block(s, "Impacto", "Menos revisión manual,\nmejor focalización.")
    bullets(s, 0.78, 2.7, 5.4, [
        "Reduce drásticamente las horas de revisión manual.",
        "Enfoca los recursos del Estado donde más se necesitan.",
        "Base para priorizar Amazonía, Orinoquía y Pacífico.",
    ])
    image_card(s, "02_mapa_cobertura", IMG_X, 1.55, IMG_W, 4.9)
    notes(s, "En conclusión, AquaBosque Minero IA reduce drásticamente las horas de revisión manual. Nos permite focalizar los recursos del Estado donde más se necesitan, con evidencia clara y unificada. Es una herramienta escalable que ya está lista para transformar datos dispersos en decisiones estratégicas. Muchas gracias.")

    # 17 · CIERRE
    s = base_slide(prs, footer=False, page=17)
    tb, tf = textbox(s, 0.75, 1.4, 7.2, 0.4)
    _run(tf.paragraphs[0], "CIERRE", 12.5, EMER, font=F_TIT, bold=True, spacing=180)
    tb, tf = textbox(s, 0.72, 2.0, 8.2, 2.2)
    _run(tf.paragraphs[0], "Datos dispersos,", 38, NAVY, font=F_TIT, bold=True)
    p = tf.add_paragraph(); _run(p, "decisiones estratégicas.", 38, EMER, font=F_TIT, bold=True)
    bullets(s, 0.78, 4.3, 7.5, [
        "Producto real, escalable y con soberanía del dato.",
        "Listo para transformar la gestión ambiental del Estado.",
    ], size=19)
    tb, tf = textbox(s, 0.78, 5.9, 6, 0.6)
    _run(tf.paragraphs[0], "Muchas gracias.", 24, FOREST, font=F_TIT, bold=True)
    for i, (name, lab) in enumerate([("qr_demo", "Demo"), ("qr_repo", "Repo")]):
        image_card(s, name, 9.8 + i * 1.7, 4.7, 1.4, 1.4, radius=1200)
        tb, tf = textbox(s, 9.8 + i * 1.7, 6.15, 1.4, 0.3)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER; _run(tf.paragraphs[0], lab, 12, GRAY, font=F_BODY)
    tb, tf = textbox(s, 0.75, 7.06, 9, 0.3)
    _run(tf.paragraphs[0], "Ministerio de Minas y Energía", 10, GRAY, font=F_BODY)
    notes(s, "En conclusión, AquaBosque Minero IA reduce drásticamente las horas de revisión manual. Nos permite focalizar los recursos del Estado donde más se necesitan, con evidencia clara y unificada. Es una herramienta escalable que ya está lista para transformar datos dispersos en decisiones estratégicas. Muchas gracias.")

    prs.save(OUT)
    print("OK:", OUT, "·", len(prs.slides._sldIdLst), "slides")


if __name__ == "__main__":
    build()
